"""
PPTX 슬라이드 생성기 (v3.1 - Impact-8 Framework + Win Theme)

[회사명] 레이어: Modern 스타일 PPTX 생성
python-pptx를 사용하여 슬라이드 생성

v3.1 추가:
- Executive Summary 슬라이드
- Next Step 슬라이드
- Differentiation 슬라이드
- Win Theme 배지 지원
"""

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .template_manager import TemplateManager
from ..schemas.proposal_schema import (
    BulletPoint,
    ContentExample,
    ChannelStrategy,
    CampaignPlan,
    KPIItem,
    ComparisonData,
    TableData,
)
from ..utils.logger import get_logger

logger = get_logger("pptx_generator")

# Modern 스타일 색상 상수
STYLE_COLORS = {
    "primary": RGBColor(0, 44, 95),        # #002C5F - 다크 블루
    "secondary": RGBColor(0, 170, 210),    # #00AAD2 - 스카이 블루
    "accent": RGBColor(230, 51, 18),       # #E63312 - 레드
    "teal": RGBColor(0, 161, 156),         # #00A19C - 청록색
    "dark_bg": RGBColor(26, 26, 26),       # #1A1A1A - 다크 배경
    "dark_blue": RGBColor(0, 44, 95),      # 별칭
    "sky_blue": RGBColor(0, 170, 210),     # 별칭
    "white": RGBColor(255, 255, 255),
    "light": RGBColor(245, 245, 245),      # #F5F5F5 - 밝은 회색
    "text_dark": RGBColor(51, 51, 51),
    "text_light": RGBColor(153, 153, 153),
    "text_gray": RGBColor(102, 102, 102),  # #666666
}


class PPTXGenerator:
    """PPTX 슬라이드 생성기"""

    def __init__(self, template_manager: TemplateManager):
        self.template_manager = template_manager
        self.design = template_manager.design_system
        self.prs: Optional[Presentation] = None

    def create_presentation(self, template_name: str = "base_template") -> Presentation:
        """새 프레젠테이션 생성"""
        self.prs = self.template_manager.load_template(template_name)
        return self.prs

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        is_part_divider: bool = False,
    ) -> None:
        """
        타이틀 슬라이드 추가

        Args:
            title: 제목
            subtitle: 부제목
            is_part_divider: Part 구분자 슬라이드 여부
        """
        layout_name = "section" if is_part_divider else "title"
        layout_idx = self.template_manager.get_layout_index(layout_name)

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[0]

        slide = self.prs.slides.add_slide(slide_layout)

        # 제목 설정
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_title_format(
                title_shape.text_frame.paragraphs[0],
                size_name="part_title" if is_part_divider else "cover_title",
            )

        # 부제목 설정
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # subtitle placeholder
                    shape.text = subtitle
                    self._apply_text_format(
                        shape.text_frame.paragraphs[0],
                        size_name="subtitle",
                        color_name="text_light",
                    )
                    break

    def add_content_slide(
        self,
        title: str,
        bullets: Optional[List[BulletPoint]] = None,
        key_message: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> None:
        """
        콘텐츠 슬라이드 추가

        Args:
            title: 슬라이드 제목
            bullets: 불릿 포인트 목록
            key_message: 핵심 메시지 (슬라이드 하단)
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("content")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[1]

        slide = self.prs.slides.add_slide(slide_layout)

        # 제목
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_title_format(slide.shapes.title.text_frame.paragraphs[0])

        # 본문 (불릿 포인트)
        if bullets:
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if body_shape:
                tf = body_shape.text_frame
                tf.word_wrap = True

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = bullet.text
                    p.level = bullet.level
                    p.font.size = self.template_manager.get_font_size("body")
                    p.font.name = self.template_manager.get_font_name("body")
                    p.font.bold = bullet.emphasis

                    if bullet.emphasis:
                        p.font.color.rgb = self.template_manager.get_color("primary")

        # 핵심 메시지 (하단)
        if key_message:
            self._add_key_message(slide, key_message)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        highlight_rows: Optional[List[int]] = None,
        notes: Optional[str] = None,
    ) -> None:
        """
        테이블 슬라이드 추가

        Args:
            title: 슬라이드 제목
            headers: 테이블 헤더
            rows: 테이블 데이터 행
            highlight_rows: 강조할 행 인덱스
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 제목 추가
        self._add_title_textbox(slide, title)

        # 테이블 생성
        rows_count = len(rows) + 1  # 헤더 포함
        cols_count = len(headers)

        if rows_count < 2 or cols_count < 1:
            return

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(12.33)
        height = Inches(min(0.5 * rows_count, 5.5))

        table = slide.shapes.add_table(
            rows_count, cols_count, left, top, width, height
        ).table

        # 열 너비 설정
        col_width = width // cols_count
        for i in range(cols_count):
            table.columns[i].width = col_width

        # 헤더 설정
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.template_manager.get_color("primary")
            self._format_table_cell(cell, is_header=True)

        # 데이터 행 설정
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= cols_count:
                    break

                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text) if cell_text else ""

                # 강조 행 처리
                if highlight_rows and row_idx in highlight_rows:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 220)

                self._format_table_cell(cell, is_header=False)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_bullets: List[BulletPoint],
        right_title: str,
        right_bullets: List[BulletPoint],
        notes: Optional[str] = None,
    ) -> None:
        """
        2단 슬라이드 추가

        Args:
            title: 슬라이드 제목
            left_title: 왼쪽 열 제목
            left_bullets: 왼쪽 열 불릿
            right_title: 오른쪽 열 제목
            right_bullets: 오른쪽 열 불릿
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 왼쪽 열
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5)
        )
        self._fill_column(left_box, left_title, left_bullets)

        # 오른쪽 열
        right_box = slide.shapes.add_textbox(
            Inches(6.93), Inches(1.5), Inches(5.9), Inches(5.5)
        )
        self._fill_column(right_box, right_title, right_bullets)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def save(self, output_path: Path) -> None:
        """프레젠테이션 저장"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"PPTX 저장 완료: {output_path}")

    def _add_title_textbox(self, slide, title: str) -> None:
        """슬라이드에 제목 텍스트박스 추가"""
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        self._apply_title_format(p)

    def _add_key_message(self, slide, message: str) -> None:
        """슬라이드 하단에 핵심 메시지 추가"""
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(12.33)
        height = Inches(0.7)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = f">> {message}"
        p.font.size = self.template_manager.get_font_size("body")
        p.font.bold = True
        p.font.color.rgb = self.template_manager.get_color("accent")
        p.alignment = PP_ALIGN.LEFT

    def _apply_title_format(self, paragraph, size_name: str = "slide_title") -> None:
        """제목 포맷 적용"""
        paragraph.font.size = self.template_manager.get_font_size(size_name)
        paragraph.font.bold = True
        paragraph.font.name = self.template_manager.get_font_name("title")
        paragraph.font.color.rgb = self.template_manager.get_color("primary")

    def _apply_text_format(
        self, paragraph, size_name: str = "body", color_name: str = "text_dark"
    ) -> None:
        """텍스트 포맷 적용"""
        paragraph.font.size = self.template_manager.get_font_size(size_name)
        paragraph.font.name = self.template_manager.get_font_name("body")
        paragraph.font.color.rgb = self.template_manager.get_color(color_name)

    def _format_table_cell(self, cell, is_header: bool = False) -> None:
        """테이블 셀 포맷 적용"""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.name = self.template_manager.get_font_name("body")
            paragraph.font.bold = is_header
            paragraph.alignment = PP_ALIGN.CENTER

            if is_header:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            else:
                paragraph.font.color.rgb = self.template_manager.get_color("text_dark")

    def _fill_column(
        self, textbox, column_title: str, bullets: List[BulletPoint]
    ) -> None:
        """2단 레이아웃의 열 채우기"""
        tf = textbox.text_frame
        tf.word_wrap = True

        # 열 제목
        p = tf.paragraphs[0]
        p.text = column_title
        p.font.size = self.template_manager.get_font_size("subtitle")
        p.font.bold = True
        p.font.color.rgb = self.template_manager.get_color("secondary")

        # 불릿 포인트
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet.text
            p.level = bullet.level
            p.font.size = self.template_manager.get_font_size("body")
            p.font.name = self.template_manager.get_font_name("body")
            p.font.bold = bullet.emphasis

    # ==========================================================================
    # 추가 레이아웃
    # ==========================================================================
    def add_three_column_slide(
        self,
        title: str,
        columns: List[dict],
        notes: Optional[str] = None,
    ) -> None:
        """
        3단 레이아웃 슬라이드 추가

        Args:
            title: 슬라이드 제목
            columns: [{"title": "열1", "content": "내용", "icon": "★"}, ...]
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 3개 열
        col_width = 3.8
        col_height = 4.5
        left_start = 0.7
        top = 1.6
        gap = 0.3

        for i, col in enumerate(columns[:3]):
            left = left_start + i * (col_width + gap)
            self._add_column_box(slide, col, left, top, col_width, col_height)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_column_box(
        self,
        slide,
        column: dict,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """3단 레이아웃의 열 박스"""
        # 배경 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        shape.line.color.rgb = self.template_manager.get_color("primary")
        shape.line.width = Pt(1.5)

        content_left = left + 0.2
        content_width = width - 0.4
        current_top = top + 0.2

        # 아이콘 (있는 경우)
        if column.get("icon"):
            icon_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top),
                Inches(content_width),
                Inches(0.6),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = column["icon"]
            icon_p.font.size = Pt(28)
            icon_p.alignment = PP_ALIGN.CENTER
            current_top += 0.6

        # 열 제목
        title_box = slide.shapes.add_textbox(
            Inches(content_left),
            Inches(current_top),
            Inches(content_width),
            Inches(0.5),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = column.get("title", "")
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = self.template_manager.get_color("primary")
        title_p.alignment = PP_ALIGN.CENTER

        # 내용
        if column.get("content"):
            content_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top + 0.6),
                Inches(content_width),
                Inches(height - current_top - 0.5),
            )
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_p = content_tf.paragraphs[0]
            content_p.text = column["content"]
            content_p.font.size = Pt(11)
            content_p.font.color.rgb = self.template_manager.get_color("text_dark")

        # 불릿 (있는 경우)
        if column.get("bullets"):
            bullets_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top + 0.6),
                Inches(content_width),
                Inches(height - current_top - 0.5),
            )
            bullets_tf = bullets_box.text_frame
            bullets_tf.word_wrap = True

            for j, bullet in enumerate(column["bullets"]):
                if j == 0:
                    p = bullets_tf.paragraphs[0]
                else:
                    p = bullets_tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(10)

    def add_big_number_slide(
        self,
        title: str,
        stats: List[dict],
        notes: Optional[str] = None,
    ) -> None:
        """
        큰 숫자 강조 슬라이드 (KPI, 성과 등)

        Args:
            title: 슬라이드 제목
            stats: [{"value": "95%", "label": "만족도", "description": "목표 대비 +5%"}, ...]
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 통계 카드들
        num_stats = min(len(stats), 4)
        total_width = 12.0
        card_width = total_width / num_stats
        left_start = 0.67
        top = 2.2

        for i, stat in enumerate(stats[:4]):
            left = left_start + i * card_width
            self._add_stat_card(slide, stat, left, top, card_width - 0.3)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_stat_card(
        self,
        slide,
        stat: dict,
        left: float,
        top: float,
        width: float,
    ) -> None:
        """통계 카드"""
        # 큰 숫자
        value_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.2),
        )
        value_tf = value_box.text_frame
        value_p = value_tf.paragraphs[0]
        value_p.text = stat.get("value", "")
        value_p.font.size = Pt(48)
        value_p.font.bold = True
        value_p.font.color.rgb = self.template_manager.get_color("primary")
        value_p.alignment = PP_ALIGN.CENTER

        # 레이블
        label_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top + 1.2),
            Inches(width),
            Inches(0.5),
        )
        label_tf = label_box.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.text = stat.get("label", "")
        label_p.font.size = Pt(14)
        label_p.font.bold = True
        label_p.font.color.rgb = self.template_manager.get_color("text_dark")
        label_p.alignment = PP_ALIGN.CENTER

        # 설명 (있는 경우)
        if stat.get("description"):
            desc_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top + 1.7),
                Inches(width),
                Inches(0.5),
            )
            desc_tf = desc_box.text_frame
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = stat["description"]
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = self.template_manager.get_color("text_light")
            desc_p.alignment = PP_ALIGN.CENTER

    def add_icon_grid_slide(
        self,
        title: str,
        items: List[dict],
        columns: int = 4,
        notes: Optional[str] = None,
    ) -> None:
        """
        아이콘 그리드 슬라이드

        Args:
            title: 슬라이드 제목
            items: [{"icon": "🎯", "title": "목표", "description": "설명"}, ...]
            columns: 열 개수 (3 또는 4)
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 그리드 아이템
        total_width = 12.0
        margin = 0.25
        item_width = (total_width - margin * (columns + 1)) / columns
        item_height = 1.8
        left_start = 0.67
        top = 1.8

        for i, item in enumerate(items[:columns * 2]):  # 최대 2행
            row = i // columns
            col = i % columns

            item_left = left_start + margin + col * (item_width + margin)
            item_top = top + row * (item_height + margin)

            self._add_icon_item(slide, item, item_left, item_top, item_width, item_height)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_icon_item(
        self,
        slide,
        item: dict,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """아이콘 그리드 아이템"""
        # 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
        shape.line.color.rgb = self.template_manager.get_color("text_light")
        shape.line.width = Pt(1)

        # 아이콘
        if item.get("icon"):
            icon_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 0.1),
                Inches(width - 0.2),
                Inches(0.5),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = item["icon"]
            icon_p.font.size = Pt(24)
            icon_p.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.6),
            Inches(width - 0.2),
            Inches(0.4),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = item.get("title", "")
        title_p.font.size = Pt(11)
        title_p.font.bold = True
        title_p.font.color.rgb = self.template_manager.get_color("primary")
        title_p.alignment = PP_ALIGN.CENTER

        # 설명
        if item.get("description"):
            desc_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 1.0),
                Inches(width - 0.2),
                Inches(height - 1.1),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = item["description"]
            desc_p.font.size = Pt(9)
            desc_p.font.color.rgb = self.template_manager.get_color("text_dark")
            desc_p.alignment = PP_ALIGN.CENTER

    def add_quote_slide(
        self,
        title: str,
        quote: str,
        author: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> None:
        """
        인용문/추천사 슬라이드

        Args:
            title: 슬라이드 제목
            quote: 인용문 내용
            author: 작성자 정보
            notes: 발표자 노트
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 큰 따옴표 배경
        quote_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5),
            Inches(2.0),
            Inches(10.33),
            Inches(3.5),
        )
        quote_bg.fill.solid()
        quote_bg.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        quote_bg.line.fill.background()

        # 인용문
        quote_box = slide.shapes.add_textbox(
            Inches(2.0),
            Inches(2.3),
            Inches(9.33),
            Inches(2.5),
        )
        quote_tf = quote_box.text_frame
        quote_tf.word_wrap = True
        quote_p = quote_tf.paragraphs[0]
        quote_p.text = f'"{quote}"'
        quote_p.font.size = Pt(20)
        quote_p.font.italic = True
        quote_p.font.color.rgb = self.template_manager.get_color("text_dark")
        quote_p.alignment = PP_ALIGN.CENTER

        # 작성자
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(6.0),
                Inches(4.8),
                Inches(5.83),
                Inches(0.5),
            )
            author_tf = author_box.text_frame
            author_p = author_tf.paragraphs[0]
            author_p.text = f"- {author}"
            author_p.font.size = Pt(14)
            author_p.font.color.rgb = self.template_manager.get_color("text_light")
            author_p.alignment = PP_ALIGN.RIGHT

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # ========================================
    # Modern 스타일 슬라이드 메서드
    # ========================================

    def add_teaser_slide(
        self,
        headline: str,
        subheadline: str = "",
        background_color: str = "dark_blue",
        notes: str = "",
    ):
        """
        티저/HOOK 슬라이드 - 임팩트 있는 오프닝
        Modern 스타일: 다크 배경, 큰 타이포그래피
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 다크 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(7.5),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS.get(background_color, STYLE_COLORS["dark_blue"])
        bg.line.fill.background()

        # 메인 헤드라인
        headline_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(2.5),
            Inches(11.33),
            Inches(2.0),
        )
        headline_tf = headline_box.text_frame
        headline_tf.word_wrap = True
        headline_p = headline_tf.paragraphs[0]
        headline_p.text = headline
        headline_p.font.size = Pt(54)
        headline_p.font.bold = True
        headline_p.font.color.rgb = RGBColor(255, 255, 255)
        headline_p.alignment = PP_ALIGN.CENTER

        # 서브 헤드라인
        if subheadline:
            sub_box = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(4.5),
                Inches(11.33),
                Inches(1.0),
            )
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = subheadline
            sub_p.font.size = Pt(24)
            sub_p.font.color.rgb = STYLE_COLORS["sky_blue"]
            sub_p.alignment = PP_ALIGN.CENTER

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_section_divider(
        self,
        phase_number: int,
        phase_title: str,
        phase_subtitle: str = "",
        notes: str = "",
    ):
        """
        섹션 구분 슬라이드 - Phase 번호와 제목
        Modern 스타일: 왼쪽에 큰 숫자, 오른쪽에 제목
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 다크 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(7.5),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        bg.line.fill.background()

        # Phase 번호 (큰 아웃라인 숫자)
        num_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(2.0),
            Inches(3.0),
            Inches(3.0),
        )
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = f"0{phase_number}" if phase_number < 10 else str(phase_number)
        num_p.font.size = Pt(120)
        num_p.font.bold = True
        num_p.font.color.rgb = STYLE_COLORS["sky_blue"]
        num_p.alignment = PP_ALIGN.CENTER

        # Phase 제목
        title_box = slide.shapes.add_textbox(
            Inches(4.5),
            Inches(2.8),
            Inches(8.0),
            Inches(1.5),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = phase_title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)

        # Phase 서브타이틀
        if phase_subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(4.5),
                Inches(4.3),
                Inches(8.0),
                Inches(1.0),
            )
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = phase_subtitle
            sub_p.font.size = Pt(20)
            sub_p.font.color.rgb = RGBColor(180, 180, 180)

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_key_message_slide(
        self,
        message: str,
        supporting_text: str = "",
        background_style: str = "gradient",
        notes: str = "",
    ):
        """
        핵심 메시지 슬라이드 - 중앙 정렬된 임팩트 메시지
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 배경 (그라데이션 효과는 단색으로 대체)
        if background_style == "dark":
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.33),
                Inches(7.5),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
            bg.line.fill.background()
            text_color = RGBColor(255, 255, 255)
        else:
            text_color = STYLE_COLORS["dark_blue"]

        # 핵심 메시지
        msg_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(2.5),
            Inches(11.33),
            Inches(2.5),
        )
        msg_tf = msg_box.text_frame
        msg_tf.word_wrap = True
        msg_p = msg_tf.paragraphs[0]
        msg_p.text = message
        msg_p.font.size = Pt(40)
        msg_p.font.bold = True
        msg_p.font.color.rgb = text_color
        msg_p.alignment = PP_ALIGN.CENTER

        # 보조 텍스트
        if supporting_text:
            sup_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(5.0),
                Inches(10.33),
                Inches(1.5),
            )
            sup_tf = sup_box.text_frame
            sup_tf.word_wrap = True
            sup_p = sup_tf.paragraphs[0]
            sup_p.text = supporting_text
            sup_p.font.size = Pt(18)
            sup_p.font.color.rgb = STYLE_COLORS["sky_blue"] if background_style == "dark" else STYLE_COLORS["text_gray"]
            sup_p.alignment = PP_ALIGN.CENTER

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_comparison_slide(
        self,
        title: str,
        as_is: dict,
        to_be: dict,
        notes: str = "",
    ):
        """
        AS-IS / TO-BE 비교 슬라이드
        Modern 스타일: 좌우 분할, 시각적 대비
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # AS-IS 영역 (왼쪽)
        as_is_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(1.5),
            Inches(6.0),
            Inches(5.5),
        )
        as_is_bg.fill.solid()
        as_is_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        as_is_bg.line.fill.background()

        # AS-IS 제목
        as_title_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.7),
            Inches(5.4),
            Inches(0.6),
        )
        as_title_tf = as_title_box.text_frame
        as_title_p = as_title_tf.paragraphs[0]
        as_title_p.text = as_is.get("title", "AS-IS (현재)")
        as_title_p.font.size = Pt(24)
        as_title_p.font.bold = True
        as_title_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # AS-IS 내용
        as_content_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(2.4),
            Inches(5.4),
            Inches(4.3),
        )
        as_content_tf = as_content_box.text_frame
        as_content_tf.word_wrap = True
        for i, item in enumerate(as_is.get("items", [])):
            p = as_content_tf.paragraphs[0] if i == 0 else as_content_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = STYLE_COLORS["text_gray"]
            p.space_after = Pt(8)

        # TO-BE 영역 (오른쪽)
        to_be_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.83),
            Inches(1.5),
            Inches(6.0),
            Inches(5.5),
        )
        to_be_bg.fill.solid()
        to_be_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        to_be_bg.line.fill.background()

        # TO-BE 제목
        to_title_box = slide.shapes.add_textbox(
            Inches(7.13),
            Inches(1.7),
            Inches(5.4),
            Inches(0.6),
        )
        to_title_tf = to_title_box.text_frame
        to_title_p = to_title_tf.paragraphs[0]
        to_title_p.text = to_be.get("title", "TO-BE (제안)")
        to_title_p.font.size = Pt(24)
        to_title_p.font.bold = True
        to_title_p.font.color.rgb = RGBColor(255, 255, 255)

        # TO-BE 내용
        to_content_box = slide.shapes.add_textbox(
            Inches(7.13),
            Inches(2.4),
            Inches(5.4),
            Inches(4.3),
        )
        to_content_tf = to_content_box.text_frame
        to_content_tf.word_wrap = True
        for i, item in enumerate(to_be.get("items", [])):
            p = to_content_tf.paragraphs[0] if i == 0 else to_content_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(8)

        # 화살표 (중앙)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(6.2),
            Inches(4.0),
            Inches(0.8),
            Inches(0.5),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = STYLE_COLORS["sky_blue"]
        arrow.line.fill.background()

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_index_slide(
        self,
        title: str = "목차",
        items: list = None,
        current_index: int = -1,
        notes: str = "",
    ):
        """
        목차/인덱스 슬라이드
        Modern 스타일: 깔끔한 번호 매기기, 현재 위치 강조
        """
        items = items or []
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 목차 항목들
        start_y = 1.8
        for i, item in enumerate(items):
            is_current = (i == current_index)

            # 번호 배경
            num_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.0),
                Inches(start_y + i * 0.8),
                Inches(0.5),
                Inches(0.5),
            )
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = STYLE_COLORS["sky_blue"] if is_current else STYLE_COLORS["dark_blue"]
            num_bg.line.fill.background()

            # 번호
            num_box = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(start_y + i * 0.8 + 0.08),
                Inches(0.5),
                Inches(0.4),
            )
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(14)
            num_p.font.bold = True
            num_p.font.color.rgb = RGBColor(255, 255, 255)
            num_p.alignment = PP_ALIGN.CENTER

            # 항목 텍스트
            item_box = slide.shapes.add_textbox(
                Inches(1.7),
                Inches(start_y + i * 0.8 + 0.05),
                Inches(10.0),
                Inches(0.5),
            )
            item_tf = item_box.text_frame
            item_p = item_tf.paragraphs[0]
            item_p.text = item
            item_p.font.size = Pt(18)
            item_p.font.bold = is_current
            item_p.font.color.rgb = STYLE_COLORS["dark_blue"] if is_current else STYLE_COLORS["text_gray"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_content_example_slide(
        self,
        title: str,
        examples: list,
        notes: str = "",
    ):
        """
        콘텐츠 예시 슬라이드 - 마케팅/PR용
        Modern 스타일: 카드형 레이아웃, 이미지 플레이스홀더
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 최대 3개 카드
        num_cards = min(len(examples), 3)
        card_width = 3.8
        gap = 0.3
        total_width = num_cards * card_width + (num_cards - 1) * gap
        start_x = (13.33 - total_width) / 2

        for i, example in enumerate(examples[:3]):
            x = start_x + i * (card_width + gap)

            # 카드 배경
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(1.6),
                Inches(card_width),
                Inches(5.5),
            )
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            card_bg.line.color.rgb = RGBColor(230, 230, 230)

            # 이미지 플레이스홀더
            img_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.15),
                Inches(1.75),
                Inches(card_width - 0.3),
                Inches(2.5),
            )
            img_placeholder.fill.solid()
            img_placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
            img_placeholder.line.fill.background()

            # 이미지 아이콘 텍스트
            icon_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(2.7),
                Inches(card_width - 0.3),
                Inches(0.5),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = "📷 이미지"
            icon_p.font.size = Pt(14)
            icon_p.font.color.rgb = RGBColor(150, 150, 150)
            icon_p.alignment = PP_ALIGN.CENTER

            # 콘텐츠 유형
            type_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(4.35),
                Inches(card_width - 0.3),
                Inches(0.4),
            )
            type_tf = type_box.text_frame
            type_p = type_tf.paragraphs[0]
            content_type = example.get("content_type", "콘텐츠")
            type_p.text = content_type
            type_p.font.size = Pt(11)
            type_p.font.bold = True
            type_p.font.color.rgb = STYLE_COLORS["sky_blue"]

            # 제목
            ex_title_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(4.7),
                Inches(card_width - 0.3),
                Inches(0.6),
            )
            ex_title_tf = ex_title_box.text_frame
            ex_title_tf.word_wrap = True
            ex_title_p = ex_title_tf.paragraphs[0]
            ex_title_p.text = example.get("title", "")
            ex_title_p.font.size = Pt(13)
            ex_title_p.font.bold = True
            ex_title_p.font.color.rgb = STYLE_COLORS["dark_blue"]

            # 설명
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(5.3),
                Inches(card_width - 0.3),
                Inches(1.0),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = example.get("description", "")[:80]
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = STYLE_COLORS["text_gray"]

            # 채널/해시태그
            channel_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(6.5),
                Inches(card_width - 0.3),
                Inches(0.4),
            )
            channel_tf = channel_box.text_frame
            channel_p = channel_tf.paragraphs[0]
            channel = example.get("channel", "")
            channel_p.text = f"#{channel}" if channel else ""
            channel_p.font.size = Pt(9)
            channel_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_channel_strategy_slide(
        self,
        title: str,
        channels: list,
        notes: str = "",
    ):
        """
        채널 전략 슬라이드 - 채널별 역할/KPI
        Modern 스타일: 채널 아이콘, 역할, KPI 표시
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 채널별 영역 (최대 4개)
        num_channels = min(len(channels), 4)
        col_width = 2.9
        gap = 0.3
        total_width = num_channels * col_width + (num_channels - 1) * gap
        start_x = (13.33 - total_width) / 2

        channel_colors = [
            STYLE_COLORS["dark_blue"],
            STYLE_COLORS["sky_blue"],
            RGBColor(230, 126, 34),  # Orange
            RGBColor(155, 89, 182),  # Purple
        ]

        for i, channel in enumerate(channels[:4]):
            x = start_x + i * (col_width + gap)
            color = channel_colors[i % len(channel_colors)]

            # 채널 헤더
            header_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(1.6),
                Inches(col_width),
                Inches(0.8),
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = color
            header_bg.line.fill.background()

            # 채널명
            name_box = slide.shapes.add_textbox(
                Inches(x),
                Inches(1.75),
                Inches(col_width),
                Inches(0.5),
            )
            name_tf = name_box.text_frame
            name_p = name_tf.paragraphs[0]
            name_p.text = channel.get("name", "")
            name_p.font.size = Pt(16)
            name_p.font.bold = True
            name_p.font.color.rgb = RGBColor(255, 255, 255)
            name_p.alignment = PP_ALIGN.CENTER

            # 역할 영역
            role_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(2.4),
                Inches(col_width),
                Inches(2.2),
            )
            role_bg.fill.solid()
            role_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            role_bg.line.color.rgb = RGBColor(230, 230, 230)

            # 역할 제목
            role_title_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(2.5),
                Inches(col_width - 0.2),
                Inches(0.4),
            )
            role_title_tf = role_title_box.text_frame
            role_title_p = role_title_tf.paragraphs[0]
            role_title_p.text = "역할"
            role_title_p.font.size = Pt(11)
            role_title_p.font.bold = True
            role_title_p.font.color.rgb = color

            # 역할 내용
            role_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(2.9),
                Inches(col_width - 0.2),
                Inches(1.5),
            )
            role_tf = role_box.text_frame
            role_tf.word_wrap = True
            role_p = role_tf.paragraphs[0]
            role_p.text = channel.get("role", "")
            role_p.font.size = Pt(10)
            role_p.font.color.rgb = STYLE_COLORS["text_gray"]

            # KPI 영역
            kpi_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(4.6),
                Inches(col_width),
                Inches(2.4),
            )
            kpi_bg.fill.solid()
            kpi_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
            kpi_bg.line.color.rgb = RGBColor(230, 230, 230)

            # KPI 제목
            kpi_title_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(4.7),
                Inches(col_width - 0.2),
                Inches(0.4),
            )
            kpi_title_tf = kpi_title_box.text_frame
            kpi_title_p = kpi_title_tf.paragraphs[0]
            kpi_title_p.text = "KPI"
            kpi_title_p.font.size = Pt(11)
            kpi_title_p.font.bold = True
            kpi_title_p.font.color.rgb = color

            # KPI 항목들
            kpis = channel.get("kpis", [])
            for j, kpi in enumerate(kpis[:3]):
                kpi_item_box = slide.shapes.add_textbox(
                    Inches(x + 0.1),
                    Inches(5.1 + j * 0.6),
                    Inches(col_width - 0.2),
                    Inches(0.55),
                )
                kpi_item_tf = kpi_item_box.text_frame
                kpi_item_tf.word_wrap = True

                # KPI 이름
                kpi_name_p = kpi_item_tf.paragraphs[0]
                kpi_name_p.text = kpi.get("name", "")
                kpi_name_p.font.size = Pt(9)
                kpi_name_p.font.color.rgb = STYLE_COLORS["text_gray"]

                # KPI 값
                kpi_value_p = kpi_item_tf.add_paragraph()
                kpi_value_p.text = kpi.get("target", "")
                kpi_value_p.font.size = Pt(12)
                kpi_value_p.font.bold = True
                kpi_value_p.font.color.rgb = color

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_campaign_slide(
        self,
        title: str,
        campaign_name: str,
        period: str,
        objective: str,
        activities: list,
        notes: str = "",
    ):
        """
        캠페인 슬라이드 - 캠페인 개요 및 활동
        Modern 스타일: 헤더 배너, 활동 타임라인
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 캠페인 헤더 배너
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.4),
            Inches(12.33),
            Inches(1.2),
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        header_bg.line.fill.background()

        # 캠페인명
        campaign_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(1.5),
            Inches(7.0),
            Inches(0.5),
        )
        campaign_tf = campaign_box.text_frame
        campaign_p = campaign_tf.paragraphs[0]
        campaign_p.text = campaign_name
        campaign_p.font.size = Pt(22)
        campaign_p.font.bold = True
        campaign_p.font.color.rgb = RGBColor(255, 255, 255)

        # 기간
        period_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(2.0),
            Inches(4.0),
            Inches(0.4),
        )
        period_tf = period_box.text_frame
        period_p = period_tf.paragraphs[0]
        period_p.text = f"📅 {period}"
        period_p.font.size = Pt(12)
        period_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 목표
        obj_box = slide.shapes.add_textbox(
            Inches(8.0),
            Inches(1.55),
            Inches(4.5),
            Inches(0.9),
        )
        obj_tf = obj_box.text_frame
        obj_tf.word_wrap = True
        obj_p = obj_tf.paragraphs[0]
        obj_p.text = objective
        obj_p.font.size = Pt(12)
        obj_p.font.color.rgb = RGBColor(200, 200, 200)
        obj_p.alignment = PP_ALIGN.RIGHT

        # 활동 타임라인/리스트
        activity_start_y = 2.9
        for i, activity in enumerate(activities[:6]):
            # 번호 원
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.7),
                Inches(activity_start_y + i * 0.75),
                Inches(0.4),
                Inches(0.4),
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = STYLE_COLORS["sky_blue"]
            num_circle.line.fill.background()

            # 번호 텍스트
            num_box = slide.shapes.add_textbox(
                Inches(0.7),
                Inches(activity_start_y + i * 0.75 + 0.05),
                Inches(0.4),
                Inches(0.35),
            )
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(12)
            num_p.font.bold = True
            num_p.font.color.rgb = RGBColor(255, 255, 255)
            num_p.alignment = PP_ALIGN.CENTER

            # 활동 내용
            act_box = slide.shapes.add_textbox(
                Inches(1.3),
                Inches(activity_start_y + i * 0.75 + 0.05),
                Inches(11.0),
                Inches(0.5),
            )
            act_tf = act_box.text_frame
            act_tf.word_wrap = True
            act_p = act_tf.paragraphs[0]

            if isinstance(activity, dict):
                act_p.text = activity.get("name", str(activity))
            else:
                act_p.text = str(activity)

            act_p.font.size = Pt(14)
            act_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_budget_slide(
        self,
        title: str,
        budget_items: list,
        total: str = "",
        notes: str = "",
    ):
        """
        예산 슬라이드 - 투자 비용 테이블
        Modern 스타일: 깔끔한 테이블, 총계 강조
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 테이블 생성
        rows = len(budget_items) + 2  # 헤더 + 항목들 + 총계
        cols = 4  # 항목, 단가, 수량, 금액

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.8),
            Inches(1.6),
            Inches(11.73),
            Inches(0.5 * rows),
        )
        table = table_shape.table

        # 헤더 스타일
        headers = ["항목", "단가", "수량", "금액"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER

        # 데이터 행
        for row_idx, item in enumerate(budget_items, start=1):
            # 항목명
            cell0 = table.cell(row_idx, 0)
            cell0.text = item.get("name", "")
            cell0.text_frame.paragraphs[0].font.size = Pt(11)

            # 단가
            cell1 = table.cell(row_idx, 1)
            cell1.text = item.get("unit_price", "")
            cell1.text_frame.paragraphs[0].font.size = Pt(11)
            cell1.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # 수량
            cell2 = table.cell(row_idx, 2)
            cell2.text = str(item.get("quantity", ""))
            cell2.text_frame.paragraphs[0].font.size = Pt(11)
            cell2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 금액
            cell3 = table.cell(row_idx, 3)
            cell3.text = item.get("amount", "")
            cell3.text_frame.paragraphs[0].font.size = Pt(11)
            cell3.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 총계 행
        total_row = rows - 1
        table.cell(total_row, 0).merge(table.cell(total_row, 2))
        total_cell = table.cell(total_row, 0)
        total_cell.text = "총계"
        total_cell.fill.solid()
        total_cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        total_cell.text_frame.paragraphs[0].font.size = Pt(12)
        total_cell.text_frame.paragraphs[0].font.bold = True
        total_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        total_amount_cell = table.cell(total_row, 3)
        total_amount_cell.text = total
        total_amount_cell.fill.solid()
        total_amount_cell.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        total_amount_cell.text_frame.paragraphs[0].font.size = Pt(14)
        total_amount_cell.text_frame.paragraphs[0].font.bold = True
        total_amount_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        total_amount_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_case_study_slide(
        self,
        title: str,
        case: dict,
        notes: str = "",
    ):
        """
        케이스 스터디 슬라이드 - 수행 실적
        Modern 스타일: 프로젝트 이미지 + 성과 KPI
        """
        layout_idx = self.template_manager.get_layout_index("blank")

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 왼쪽: 이미지 플레이스홀더 + 프로젝트 정보
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.5),
            Inches(6.5),
            Inches(4.0),
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        img_placeholder.line.fill.background()

        # 이미지 아이콘
        icon_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(3.2),
            Inches(6.5),
            Inches(0.5),
        )
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = "📷 프로젝트 이미지"
        icon_p.font.size = Pt(16)
        icon_p.font.color.rgb = RGBColor(150, 150, 150)
        icon_p.alignment = PP_ALIGN.CENTER

        # 프로젝트명
        proj_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(5.6),
            Inches(6.5),
            Inches(0.5),
        )
        proj_tf = proj_box.text_frame
        proj_p = proj_tf.paragraphs[0]
        proj_p.text = case.get("project_name", "")
        proj_p.font.size = Pt(18)
        proj_p.font.bold = True
        proj_p.font.color.rgb = STYLE_COLORS["dark_blue"]

        # 클라이언트/기간
        info_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6.1),
            Inches(6.5),
            Inches(0.8),
        )
        info_tf = info_box.text_frame
        info_p = info_tf.paragraphs[0]
        client = case.get("client", "")
        period = case.get("period", "")
        info_p.text = f"{client} | {period}"
        info_p.font.size = Pt(12)
        info_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 오른쪽: 성과 KPI 영역
        kpi_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2),
            Inches(1.5),
            Inches(5.63),
            Inches(5.5),
        )
        kpi_bg.fill.solid()
        kpi_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        kpi_bg.line.fill.background()

        # 성과 제목
        result_title_box = slide.shapes.add_textbox(
            Inches(7.4),
            Inches(1.7),
            Inches(5.23),
            Inches(0.5),
        )
        result_title_tf = result_title_box.text_frame
        result_title_p = result_title_tf.paragraphs[0]
        result_title_p.text = "주요 성과"
        result_title_p.font.size = Pt(16)
        result_title_p.font.bold = True
        result_title_p.font.color.rgb = RGBColor(255, 255, 255)

        # KPI 항목들
        kpis = case.get("kpis", case.get("results", []))
        for i, kpi in enumerate(kpis[:4]):
            kpi_box = slide.shapes.add_textbox(
                Inches(7.4),
                Inches(2.4 + i * 1.2),
                Inches(5.23),
                Inches(1.0),
            )
            kpi_tf = kpi_box.text_frame

            # KPI 값 (큰 숫자)
            kpi_value_p = kpi_tf.paragraphs[0]
            if isinstance(kpi, dict):
                kpi_value_p.text = kpi.get("value", kpi.get("target", ""))
            else:
                kpi_value_p.text = str(kpi)
            kpi_value_p.font.size = Pt(28)
            kpi_value_p.font.bold = True
            kpi_value_p.font.color.rgb = STYLE_COLORS["sky_blue"]

            # KPI 이름
            kpi_name_p = kpi_tf.add_paragraph()
            if isinstance(kpi, dict):
                kpi_name_p.text = kpi.get("name", kpi.get("metric", ""))
            else:
                kpi_name_p.text = ""
            kpi_name_p.font.size = Pt(11)
            kpi_name_p.font.color.rgb = RGBColor(200, 200, 200)

        # 프로젝트 설명 (아래)
        desc = case.get("description", case.get("overview", ""))
        if desc:
            desc_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(6.7),
                Inches(12.33),
                Inches(0.6),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = desc[:150] + "..." if len(desc) > 150 else desc
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # ========================================
    # v3.1 추가: Executive Summary, Next Step, Differentiation
    # ========================================

    def add_executive_summary_slide(
        self,
        project_objective: str,
        win_themes: list,
        kpis: list,
        why_us_points: list,
        notes: str = "",
    ):
        """
        Executive Summary 슬라이드 - 의사결정권자용 1페이지 요약

        Args:
            project_objective: 프로젝트 핵심 목표 (1문장)
            win_themes: [{"name": "Win Theme명", "description": "설명"}, ...]
            kpis: [{"metric": "KPI명", "target": "목표값", "basis": "산출근거"}, ...]
            why_us_points: ["포인트1", "포인트2", ...]
        """
        layout_idx = self.template_manager.get_layout_index("blank")
        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 왼쪽 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = STYLE_COLORS["primary"]
        accent_bar.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.6))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = "EXECUTIVE SUMMARY"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["primary"]

        # 프로젝트 목표
        obj_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6)
        )
        obj_bg.fill.solid()
        obj_bg.fill.fore_color.rgb = STYLE_COLORS["primary"]
        obj_bg.line.fill.background()

        obj_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.4))
        obj_p = obj_box.text_frame.paragraphs[0]
        obj_p.text = project_objective
        obj_p.font.size = Pt(16)
        obj_p.font.bold = True
        obj_p.font.color.rgb = STYLE_COLORS["white"]

        # Win Themes (3개 카드)
        win_colors = [STYLE_COLORS["primary"], STYLE_COLORS["secondary"], STYLE_COLORS["teal"]]
        for i, theme in enumerate(win_themes[:3]):
            x = Inches(0.5 + i * 4.2)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(4.0), Inches(1.4)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = win_colors[i % 3]
            card.line.fill.background()

            name_box = slide.shapes.add_textbox(x, Inches(2.1), Inches(4.0), Inches(0.5))
            name_p = name_box.text_frame.paragraphs[0]
            name_p.text = theme.get("name", "")
            name_p.font.size = Pt(14)
            name_p.font.bold = True
            name_p.font.color.rgb = STYLE_COLORS["white"]
            name_p.alignment = PP_ALIGN.CENTER

            desc_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.6), Inches(3.8), Inches(0.7))
            desc_box.text_frame.word_wrap = True
            desc_p = desc_box.text_frame.paragraphs[0]
            desc_p.text = theme.get("description", "")
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = STYLE_COLORS["white"]
            desc_p.alignment = PP_ALIGN.CENTER

        # KPI 카드 (4개)
        for i, kpi in enumerate(kpis[:4]):
            x = Inches(0.5 + i * 3.2)
            kpi_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.6), Inches(3.0), Inches(1.5)
            )
            kpi_card.fill.solid()
            kpi_card.fill.fore_color.rgb = STYLE_COLORS["light"]
            kpi_card.line.fill.background()

            metric_box = slide.shapes.add_textbox(x, Inches(3.7), Inches(3.0), Inches(0.35))
            metric_p = metric_box.text_frame.paragraphs[0]
            metric_p.text = kpi.get("metric", "")
            metric_p.font.size = Pt(14)
            metric_p.font.bold = True
            metric_p.font.color.rgb = STYLE_COLORS["primary"]
            metric_p.alignment = PP_ALIGN.CENTER

            target_box = slide.shapes.add_textbox(x, Inches(4.05), Inches(3.0), Inches(0.4))
            target_p = target_box.text_frame.paragraphs[0]
            target_p.text = kpi.get("target", "")
            target_p.font.size = Pt(18)
            target_p.font.bold = True
            target_p.font.color.rgb = STYLE_COLORS["text_dark"]
            target_p.alignment = PP_ALIGN.CENTER

            basis_box = slide.shapes.add_textbox(x, Inches(4.5), Inches(3.0), Inches(0.55))
            basis_box.text_frame.word_wrap = True
            basis_p = basis_box.text_frame.paragraphs[0]
            basis_p.text = kpi.get("basis", kpi.get("calculation_basis", ""))
            basis_p.font.size = Pt(9)
            basis_p.font.color.rgb = STYLE_COLORS["text_gray"]
            basis_p.alignment = PP_ALIGN.CENTER

        # Why Us
        why_text = "  ".join([f"✓ {p}" for p in why_us_points[:4]])
        why_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5))
        why_box.text_frame.word_wrap = True
        why_p = why_box.text_frame.paragraphs[0]
        why_p.text = why_text
        why_p.font.size = Pt(12)
        why_p.font.bold = True
        why_p.font.color.rgb = STYLE_COLORS["secondary"]
        why_p.alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_next_step_slide(
        self,
        headline: str,
        steps: list,
        call_to_action: list,
        contact_info: dict = None,
        notes: str = "",
    ):
        """
        Next Step 슬라이드 - 다음 단계 안내 / Call to Action

        Args:
            headline: 헤드라인
            steps: [{"title": "계약 체결", "date": "2026.03", "description": "계약 협의"}, ...]
            call_to_action: ["10개월간 브랜드 인지도 +20%p 달성", ...]
            contact_info: {"name": "담당자명", "phone": "전화번호", "email": "이메일"}
        """
        layout_idx = self.template_manager.get_layout_index("blank")
        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 왼쪽 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = STYLE_COLORS["primary"]
        accent_bar.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.6))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = "NEXT STEP"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["primary"]

        # 헤드라인
        headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
        headline_p = headline_box.text_frame.paragraphs[0]
        headline_p.text = headline
        headline_p.font.size = Pt(24)
        headline_p.font.bold = True
        headline_p.font.color.rgb = STYLE_COLORS["text_dark"]

        # Step 카드들
        step_colors = [STYLE_COLORS["primary"]] + [STYLE_COLORS["secondary"]] * 10
        for i, step in enumerate(steps[:4]):
            x = Inches(0.5 + i * 3.2)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.0), Inches(1.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = step_colors[min(i, len(step_colors) - 1)]
            card.line.fill.background()

            num_box = slide.shapes.add_textbox(x, Inches(1.9), Inches(3.0), Inches(0.3))
            num_p = num_box.text_frame.paragraphs[0]
            num_p.text = f"STEP {i + 1}"
            num_p.font.size = Pt(11)
            num_p.font.bold = True
            num_p.font.color.rgb = STYLE_COLORS["white"]
            num_p.alignment = PP_ALIGN.CENTER

            title_b = slide.shapes.add_textbox(x, Inches(2.2), Inches(3.0), Inches(0.4))
            title_p2 = title_b.text_frame.paragraphs[0]
            title_p2.text = step.get("title", "")
            title_p2.font.size = Pt(18)
            title_p2.font.bold = True
            title_p2.font.color.rgb = STYLE_COLORS["white"]
            title_p2.alignment = PP_ALIGN.CENTER

            date_box = slide.shapes.add_textbox(x, Inches(2.6), Inches(3.0), Inches(0.3))
            date_p = date_box.text_frame.paragraphs[0]
            date_p.text = step.get("date", "")
            date_p.font.size = Pt(12)
            date_p.font.color.rgb = STYLE_COLORS["white"]
            date_p.alignment = PP_ALIGN.CENTER

            desc_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.95), Inches(2.8), Inches(0.55))
            desc_box.text_frame.word_wrap = True
            desc_p = desc_box.text_frame.paragraphs[0]
            desc_p.text = step.get("description", "")
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = STYLE_COLORS["white"]
            desc_p.alignment = PP_ALIGN.CENTER

            if i < len(steps) - 1:
                arrow_box = slide.shapes.add_textbox(x + Inches(3.05), Inches(2.5), Inches(0.3), Inches(0.4))
                arrow_p = arrow_box.text_frame.paragraphs[0]
                arrow_p.text = "→"
                arrow_p.font.size = Pt(20)
                arrow_p.font.bold = True
                arrow_p.font.color.rgb = STYLE_COLORS["text_light"]

        # CTA 영역
        cta_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.2)
        )
        cta_bg.fill.solid()
        cta_bg.fill.fore_color.rgb = STYLE_COLORS["light"]
        cta_bg.line.fill.background()

        cta_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.35))
        cta_title_p = cta_title.text_frame.paragraphs[0]
        cta_title_p.text = "저희가 제안하는 것"
        cta_title_p.font.size = Pt(14)
        cta_title_p.font.bold = True
        cta_title_p.font.color.rgb = STYLE_COLORS["primary"]

        for i, cta in enumerate(call_to_action[:4]):
            x = Inches(0.9) if i < 2 else Inches(6.5)
            y = Inches(4.4 + (i % 2) * 0.35)
            cta_box = slide.shapes.add_textbox(x, y, Inches(5.4), Inches(0.35))
            cta_p = cta_box.text_frame.paragraphs[0]
            cta_p.text = f"✓ {cta}"
            cta_p.font.size = Pt(12)
            cta_p.font.color.rgb = STYLE_COLORS["text_dark"]

        # 연락처
        if contact_info:
            contact_text = f"담당자: {contact_info.get('name', '[담당자명]')} | 연락처: {contact_info.get('phone', '[전화번호]')} | 이메일: {contact_info.get('email', '[이메일]')}"
        else:
            contact_text = "담당자: [담당자명] | 연락처: [전화번호] | 이메일: [이메일]"

        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4))
        contact_p = contact_box.text_frame.paragraphs[0]
        contact_p.text = contact_text
        contact_p.font.size = Pt(12)
        contact_p.font.color.rgb = STYLE_COLORS["text_gray"]
        contact_p.alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_section_divider_with_win_theme(
        self,
        phase_number: int,
        phase_title: str,
        phase_subtitle: str = "",
        story_title: str = "",
        win_theme: str = "",
        notes: str = "",
    ):
        """
        섹션 구분 슬라이드 (Win Theme 배지 포함)
        """
        layout_idx = self.template_manager.get_layout_index("blank")
        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.prs.slide_layouts[6]

        slide = self.prs.slides.add_slide(slide_layout)

        # 다크 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS["dark_bg"]
        bg.line.fill.background()

        # Phase 번호
        num_text = f"0{phase_number}" if phase_number < 10 else str(phase_number)
        num_box = slide.shapes.add_textbox(Inches(5.0), Inches(0.5), Inches(8), Inches(6))
        num_p = num_box.text_frame.paragraphs[0]
        num_p.text = num_text
        num_p.font.size = Pt(320)
        num_p.font.bold = True
        num_p.font.color.rgb = RGBColor(40, 40, 40)
        num_p.alignment = PP_ALIGN.RIGHT

        # Part 라벨
        part_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(4), Inches(0.5))
        part_p = part_box.text_frame.paragraphs[0]
        part_p.text = f"PART {num_text}"
        part_p.font.size = Pt(16)
        part_p.font.bold = True
        part_p.font.color.rgb = STYLE_COLORS["secondary"]

        # 스토리 타이틀
        if story_title:
            story_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(8), Inches(0.5))
            story_p = story_box.text_frame.paragraphs[0]
            story_p.text = story_title
            story_p.font.size = Pt(20)
            story_p.font.color.rgb = STYLE_COLORS["secondary"]

        # 메인 타이틀
        y_title = Inches(3.9) if story_title else Inches(3.6)
        title_box = slide.shapes.add_textbox(Inches(0.8), y_title, Inches(8), Inches(1.0))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = phase_title
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["white"]

        # 서브타이틀
        if phase_subtitle:
            y_sub = Inches(5.3) if story_title else Inches(5.0)
            sub_box = slide.shapes.add_textbox(Inches(0.8), y_sub, Inches(8), Inches(0.5))
            sub_p = sub_box.text_frame.paragraphs[0]
            sub_p.text = phase_subtitle
            sub_p.font.size = Pt(16)
            sub_p.font.color.rgb = STYLE_COLORS["text_light"]

        # Win Theme 배지
        if win_theme:
            badge_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(8), Inches(0.6)
            )
            badge_bg.fill.solid()
            badge_bg.fill.fore_color.rgb = STYLE_COLORS["secondary"]
            badge_bg.line.fill.background()

            badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.58), Inches(8), Inches(0.45))
            badge_p = badge_box.text_frame.paragraphs[0]
            badge_p.text = f"💡 Win Theme: {win_theme}"
            badge_p.font.size = Pt(14)
            badge_p.font.bold = True
            badge_p.font.color.rgb = STYLE_COLORS["white"]

        if notes:
            slide.notes_slide.notes_text_frame.text = notes
