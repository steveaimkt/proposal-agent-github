#!/usr/bin/env python3
"""
slide_kit.py — 입찰 제안서 PPTX 공통 렌더링 툴킷 v3.6

컨설팅 스타일 + Modern 컬러 시스템
- 미세 그림자 + 그라디언트 깊이감
- 라운드 코너 카드 + 직각 도형 혼용
- 도형 위 텍스트는 항상 중앙 정렬
- 매 페이지 도식화/구조 중심
- Action Title (인사이트 기반 제목)
- Source 출처 표기
- 1 슬라이드 = 1 인사이트

v3.6: 컬러 유틸(darken/lighten) + 그라디언트 커버/클로징 + 그림자 확대 적용
      + SemiBold 타이포 + KPIS/GRID/TABLE 폴리시 + LINE_CHART smooth 버그 수정

사용법:
    import importlib.util
    spec = importlib.util.spec_from_file_location('sk', '경로/src/generators/slide_kit.py')
    sk = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(sk)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import os


# ═══════════════════════════════════════════════════════════════
#  1. 디자인 상수
# ═══════════════════════════════════════════════════════════════

# ── 컬러 유틸리티 (v3.6) ──────────────────────────────────────

def _darken(r, g, b, amount=0.25):
    """RGB를 amount 비율만큼 어둡게"""
    return (max(0, int(r * (1 - amount))),
            max(0, int(g * (1 - amount))),
            max(0, int(b * (1 - amount))))


def _lighten(r, g, b, amount=0.35):
    """RGB를 amount 비율만큼 밝게"""
    return (min(255, int(r + (255 - r) * amount)),
            min(255, int(g + (255 - g) * amount)),
            min(255, int(b + (255 - b) * amount)))


def darken(color, amount=0.25):
    """RGBColor를 amount만큼 어둡게 → 새 RGBColor 반환"""
    r, g, b = color[0], color[1], color[2]
    dr, dg, db = _darken(r, g, b, amount)
    return RGBColor(dr, dg, db)


def lighten(color, amount=0.35):
    """RGBColor를 amount만큼 밝게 → 새 RGBColor 반환"""
    r, g, b = color[0], color[1], color[2]
    lr, lg, lb = _lighten(r, g, b, amount)
    return RGBColor(lr, lg, lb)


# 컬러 팔레트 (Modern 절제 스타일)
C = {
    "primary":   RGBColor(0, 44, 95),       # #002C5F  주색 (다크블루)
    "secondary": RGBColor(0, 170, 210),      # #00AAD2  보조색 (스카이블루)
    "teal":      RGBColor(0, 161, 156),      # #00A19C  틸
    "accent":    RGBColor(230, 51, 18),      # #E63312  강조 (레드)
    "dark":      RGBColor(33, 33, 33),       # #212121  본문 기본색
    "light":     RGBColor(245, 245, 245),    # #F5F5F5  밝은 배경
    "white":     RGBColor(255, 255, 255),    # #FFFFFF
    "gray":      RGBColor(117, 117, 117),    # #757575  보조 텍스트
    "lgray":     RGBColor(200, 200, 200),    # #C8C8C8  구분선
    "green":     RGBColor(46, 125, 50),      # #2E7D32  성과/긍정
    "orange":    RGBColor(245, 166, 35),     # #F5A623  주의
    "gold":      RGBColor(197, 151, 62),     # #C5973E  프리미엄
}

# 파생 컬러 (v3.6) — 깊이감/계층 표현용
C["primary_dark"]   = darken(C["primary"], 0.3)      # 더 진한 네이비
C["primary_light"]  = lighten(C["primary"], 0.85)     # 연한 블루 배경
C["secondary_dark"] = darken(C["secondary"], 0.25)    # 진한 스카이블루
C["secondary_light"]= lighten(C["secondary"], 0.80)   # 연한 스카이 배경
C["teal_light"]     = lighten(C["teal"], 0.80)        # 연한 틸 배경
C["accent_light"]   = lighten(C["accent"], 0.80)      # 연한 레드 배경
C["green_light"]    = lighten(C["green"], 0.80)       # 연한 그린 배경
C["card_bg"]        = RGBColor(250, 250, 252)         # 카드 배경 (약간 블루)
C["card_border"]    = RGBColor(230, 232, 236)         # 카드 테두리

# 슬라이드 규격 (16:9)
SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.8)       # 좌측 여백
MR = Inches(0.8)       # 우측 여백
MT_Y = Inches(0.4)     # 상단 여백
CW = SW - ML - MR      # 콘텐츠 너비

# 타이포그래피
FONT = "Pretendard"

# 폰트 웨이트 (v3.6) — python-pptx는 font.name으로 웨이트 구분
FONT_W = {
    "light":    "Pretendard Light",
    "regular":  "Pretendard",
    "medium":   "Pretendard Medium",
    "semibold": "Pretendard SemiBold",
    "bold":     "Pretendard Bold",
    "black":    "Pretendard Black",
}

# 폰트 사이즈 체계
SZ = {
    "hero":      60,   # 표지
    "divider":   40,   # 섹션 구분자
    "action":    20,   # Action Title (슬라이드 제목)
    "subtitle":  16,   # 부제
    "body":      13,   # 본문
    "body_sm":   11,   # 본문 소
    "caption":   10,   # 캡션
    "source":     8,   # 출처
}


# ═══════════════════════════════════════════════════════════════
#  2. 프레젠테이션 / 슬라이드 생성
# ═══════════════════════════════════════════════════════════════

def new_presentation():
    """16:9 빈 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def new_slide(prs):
    """빈 레이아웃 슬라이드 추가 (흰 배경)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["white"]
    return s


# ═══════════════════════════════════════════════════════════════
#  3. 기본 도형
# ═══════════════════════════════════════════════════════════════

def R(s, l, t, w, h, f=None, lc=None, lw=1):
    """직각 사각형

    Args:
        f: 채우기 색상, None이면 투명
        lc: 테두리 색상, None이면 없음
        lw: 테두리 두께(pt)
    """
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if f:
        sh.fill.solid()
        sh.fill.fore_color.rgb = f
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc
        sh.line.width = Pt(lw)
    return sh


def BOX(s, l, t, w, h, f, text="", sz=13, tc=None, b=False):
    """텍스트가 중앙 정렬된 직각 박스 (도형 + 텍스트 일체형)

    도형 위에 텍스트가 있을 때 항상 이 함수 사용.
    텍스트는 수평/수직 모두 중앙 정렬됨.
    """
    if tc is None:
        tc = C["white"]
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = f
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def OBOX(s, l, t, w, h, text="", sz=13, tc=None, b=False, lc=None):
    """테두리만 있는 박스 (아웃라인 박스) — 텍스트 중앙 정렬

    배경 투명, 테두리만 표시.
    """
    if tc is None:
        tc = C["dark"]
    if lc is None:
        lc = C["primary"]
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    sh.line.color.rgb = lc
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


# ═══════════════════════════════════════════════════════════════
#  4. 텍스트 헬퍼
# ═══════════════════════════════════════════════════════════════

def T(s, l, t, w, h, text, sz=13, c=None, b=False, al=PP_ALIGN.LEFT, ls=1.4,
      fn=None):
    """단일 스타일 텍스트

    Args:
        fn: 폰트 이름 (None이면 FONT 사용, FONT_W["semibold"] 등 사용 가능)
    """
    if c is None:
        c = C["dark"]
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(0))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn or FONT
    p.alignment = al
    p.line_spacing = Pt(int(sz * ls))
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def RT(s, l, t, w, h, parts, al=PP_ALIGN.LEFT, ls=1.4):
    """리치 텍스트 — [(text, size, color, bold), ...]"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(0))
    p = tf.paragraphs[0]
    p.alignment = al
    max_sz = 13
    for text, sz, c, b in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = c
        r.font.bold = b
        r.font.name = FONT
        if sz > max_sz:
            max_sz = sz
    p.line_spacing = Pt(int(max_sz * ls))
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def MT(s, l, t, w, h, lines, sz=13, c=None, b=False, al=PP_ALIGN.LEFT, ls=1.6, bul=False):
    """멀티라인 텍스트"""
    if c is None:
        c = C["dark"]
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + ln) if bul else ln
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = FONT
        p.alignment = al
        p.line_spacing = Pt(int(sz * ls))
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tb


# ═══════════════════════════════════════════════════════════════
#  5. 슬라이드 공통 요소
# ═══════════════════════════════════════════════════════════════

def bg(s, c):
    """단색 배경"""
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def gradient_bg(s, c1, c2, angle=270.0):
    """그래디언트 배경"""
    fill = s.background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = angle


def gradient_shape(shape, c1, c2, angle=270.0):
    """도형에 그라디언트 채우기 적용 (v3.6)"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = angle
    return shape


# 그라디언트 프리셋 (v3.6)
GRAD = {
    "cover":     lambda: (C["primary_dark"], C["primary"]),     # 커버용 다크→블루
    "closing":   lambda: (C["primary"], darken(C["primary"], 0.15)),  # 클로징
    "highlight": lambda: (C["primary"], C["secondary_dark"]),   # 하이라이트 박스
    "section":   lambda: (darken(C["dark"], 0.3), C["dark"]),   # 섹션 구분자
    "teal":      lambda: (darken(C["teal"], 0.2), C["teal"]),   # 틸 계열
    "accent":    lambda: (C["accent"], lighten(C["accent"], 0.2)),  # 강조 계열
}


def set_char_spacing(tb, spacing=200):
    """자간 설정 (100 = 1pt)"""
    try:
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(spacing))
    except Exception:
        pass


def PN(s, n):
    """페이지 번호 (우하단)"""
    T(s, SW - Inches(1.0), SH - Inches(0.4), Inches(0.7), Inches(0.25),
      str(n), sz=SZ["source"], c=C["gray"], al=PP_ALIGN.RIGHT)


def SRC(s, text):
    """출처 표기 (좌하단) — McKinsey 스타일"""
    T(s, ML, SH - Inches(0.4), Inches(8), Inches(0.25),
      f"Source: {text}", sz=SZ["source"], c=C["gray"])


def TB(s, text, pg=None, src=None):
    """Action Title 상단바 (McKinsey 스타일)

    - 좌측 프라이머리 라인
    - Action Title = 인사이트 기반 제목 (문장형)
    - 하단 구분선
    """
    R(s, Inches(0), Inches(0), Inches(0.08), SH, f=C["primary"])
    T(s, ML, Inches(0.35), CW, Inches(0.55),
      text, sz=SZ["action"], c=C["dark"], b=True,
      fn=FONT_W["semibold"])
    R(s, ML, Inches(0.88), CW, Pt(1), f=C["lgray"])
    if pg:
        PN(s, pg)
    if src:
        SRC(s, src)


def WB(s, theme_key, win_themes, x=None, y=None, w=None):
    """Win Theme 뱃지"""
    if x is None:
        x = ML
    if y is None:
        y = SH - Inches(0.9)
    if w is None:
        w = Inches(4.5)
    BOX(s, x, y, w, Inches(0.35), C["teal"],
        f"Win Theme  |  {win_themes.get(theme_key, theme_key)}",
        sz=SZ["caption"], tc=C["white"], b=True)


def IMG(s, l, t, w, h, desc="이미지 영역"):
    """이미지 플레이스홀더"""
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = C["light"]
    sh.line.color.rgb = C["lgray"]
    sh.line.width = Pt(1)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(SZ["caption"])
    r.font.color.rgb = C["gray"]
    r.font.name = FONT
    return sh


# ═══════════════════════════════════════════════════════════════
#  6. 도식화 헬퍼 — 구조를 잡는 함수들
# ═══════════════════════════════════════════════════════════════

def FLOW(s, items, y=None, h=None, colors=None):
    """프로세스 플로우 (가로 화살표 연결)

    Args:
        items: [("제목", "설명"), ...] 리스트 (3~5개 권장)
        y: 시작 Y 위치 (기본 1.2")
        h: 박스 높이 (기본 1.2")
        colors: 색상 리스트, None이면 primary 계열 자동
    """
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.2)
    n = len(items)
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
        colors = [palette[i % len(palette)] for i in range(n)]
    arrow_w = 0.35
    total = float(CW / 914400)
    box_w = (total - arrow_w * (n - 1)) / n
    for i, (title, desc) in enumerate(items):
        x = ML + Inches((box_w + arrow_w) * i)
        BOX(s, x, y, Inches(box_w), h, colors[i],
            title, sz=SZ["body"], tc=C["white"], b=True)
        if desc:
            T(s, x, y + h + Inches(0.1), Inches(box_w), Inches(0.8),
              desc, sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.5)
        if i < n - 1:
            ax = x + Inches(box_w)
            T(s, ax, y, Inches(arrow_w), h,
              "→", sz=20, c=C["gray"], b=True, al=PP_ALIGN.CENTER)


def COLS(s, items, y=None, h=None, colors=None, show_header=True, shadow=True):
    """N-컬럼 카드 레이아웃 (v3.6 — 그림자 + 시각 폴리시)

    Args:
        items: [{"title": "제목", "body": ["항목1", "항목2"]}, ...] (2~4개)
        y: 시작 Y (기본 1.2")
        h: 카드 높이 (기본 3.5")
        colors: 헤더 색상 리스트
        show_header: False면 헤더 없이 아웃라인 박스
        shadow: 그림자 적용 여부
    """
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(3.5)
    n = len(items)
    gap = 0.2
    total = float(CW / 914400)
    col_w = (total - gap * (n - 1)) / n
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
        colors = [palette[i % len(palette)] for i in range(n)]
    header_h = Inches(0.5)
    for i, item in enumerate(items):
        x = ML + Inches((col_w + gap) * i)
        title = item.get("title", "")
        body = item.get("body", [])
        if show_header:
            # 카드 배경 (그림자)
            card_sh = R(s, x, y, Inches(col_w), h,
                        f=C["card_bg"], lc=C["card_border"], lw=0.5)
            if shadow:
                add_shadow(card_sh, preset="card")
            BOX(s, x, y, Inches(col_w), header_h, colors[i],
                title, sz=SZ["body"], tc=C["white"], b=True)
            MT(s, x + Inches(0.15), y + header_h + Inches(0.1),
               Inches(col_w - 0.3), h - header_h - Inches(0.2),
               body, sz=SZ["body_sm"], bul=True)
        else:
            header_h_ns = Inches(0.5)
            OBOX(s, x, y, Inches(col_w), header_h_ns, title, sz=SZ["body"],
                 tc=C["primary"], b=True, lc=colors[i])
            R(s, x, y + header_h_ns, Inches(col_w), h - header_h_ns,
              lc=colors[i], lw=1.5)
            MT(s, x + Inches(0.15), y + header_h_ns + Inches(0.1),
               Inches(col_w - 0.3), h - header_h_ns - Inches(0.2),
               body, sz=SZ["body_sm"], bul=True)


def PYRAMID(s, levels, y=None, w_max=None, h_total=None):
    """피라미드 구조 (위가 좁고 아래가 넓음 — McKinsey Pyramid Principle)

    Args:
        levels: [("최상위 메시지", color), ("중간", color), ("하단", color)]
                위에서 아래 순서
        y: 시작 Y (기본 1.2")
        w_max: 최대 너비 (기본 CW)
        h_total: 전체 높이 (기본 4.5")
    """
    if y is None:
        y = Inches(1.2)
    if w_max is None:
        w_max = float(CW / 914400)
    if h_total is None:
        h_total = 4.5
    n = len(levels)
    level_h = h_total / n
    center = float(ML / 914400) + w_max / 2
    for i, (text, clr) in enumerate(levels):
        ratio = 0.4 + 0.6 * (i / max(n - 1, 1))
        lw = w_max * ratio
        lx = center - lw / 2
        ly = float(y / 914400) + level_h * i
        BOX(s, Inches(lx), Inches(ly), Inches(lw), Inches(level_h - 0.08),
            clr, text, sz=SZ["body"], tc=C["white"], b=True)


def MATRIX(s, quadrants, x_label="", y_label="", y_start=None):
    """2x2 매트릭스 (McKinsey 전략 매트릭스)

    Args:
        quadrants: [("좌상", color), ("우상", color), ("좌하", color), ("우하", color)]
        x_label: X축 라벨
        y_label: Y축 라벨
        y_start: 시작 Y
    """
    if y_start is None:
        y_start = Inches(1.2)
    total = float(CW / 914400)
    label_w = 0.5
    gap = 0.1
    cell_w = (total - label_w - gap * 2) / 2
    cell_h = 2.2
    ox = float(ML / 914400) + label_w + gap
    oy = float(y_start / 914400)
    # Y축 라벨
    if y_label:
        T(s, ML, y_start, Inches(label_w), Inches(cell_h * 2 + gap),
          y_label, sz=SZ["body_sm"], c=C["gray"], b=True,
          al=PP_ALIGN.CENTER)
    # X축 라벨
    if x_label:
        T(s, Inches(ox), Inches(oy + cell_h * 2 + gap + 0.1),
          Inches(cell_w * 2 + gap), Inches(0.3),
          x_label, sz=SZ["body_sm"], c=C["gray"], b=True,
          al=PP_ALIGN.CENTER)
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for idx, (text, clr) in enumerate(quadrants[:4]):
        col, row = positions[idx]
        bx = Inches(ox + col * (cell_w + gap))
        by = Inches(oy + row * (cell_h + gap))
        BOX(s, bx, by, Inches(cell_w), Inches(cell_h), clr,
            text, sz=SZ["body"], tc=C["white"], b=True)


def TABLE(s, headers, rows, y=None, col_widths=None):
    """데이터 테이블 (McKinsey 비교표 스타일)

    Args:
        headers: ["항목", "AS-IS", "TO-BE"]
        rows: [["항목1", "현재", "목표"], ...]
        y: 시작 Y
        col_widths: 열 너비 비율 리스트, None이면 균등
    """
    if y is None:
        y = Inches(1.2)
    n_cols = len(headers)
    total = float(CW / 914400)
    if col_widths is None:
        widths = [total / n_cols] * n_cols
    else:
        ratio_sum = sum(col_widths)
        widths = [total * (r / ratio_sum) for r in col_widths]
    row_h = 0.45
    # 헤더 (v3.6 — darken primary + SemiBold)
    cx = float(ML / 914400)
    for j, hdr in enumerate(headers):
        BOX(s, Inches(cx), y, Inches(widths[j]), Inches(row_h),
            C["primary"], hdr, sz=SZ["body_sm"], tc=C["white"], b=True)
        cx += widths[j]
    # 데이터 행 (v3.6 — 미세 하단 구분선 추가)
    for i, row in enumerate(rows):
        cx = float(ML / 914400)
        ry = float(y / 914400) + row_h * (i + 1)
        bgc = C["card_bg"] if i % 2 == 0 else C["white"]
        for j, cell in enumerate(row):
            BOX(s, Inches(cx), Inches(ry), Inches(widths[j]), Inches(row_h),
                bgc, cell, sz=SZ["body_sm"], tc=C["dark"])
            cx += widths[j]
        # 행 하단 구분선
        R(s, ML, Inches(ry + row_h), CW, Pt(0.5), f=C["card_border"])


def HIGHLIGHT(s, text, sub="", y=None, color=None, grad=False):
    """핵심 메시지 강조 박스 (v3.6 — 그라디언트 옵션 + 라운드)

    Args:
        text: 메인 메시지
        sub: 보조 텍스트
        y: Y 위치
        color: 배경색
        grad: True면 그라디언트 적용 (color→secondary_dark)
    """
    if y is None:
        y = Inches(1.2)
    if color is None:
        color = C["primary"]
    if sub:
        h = Inches(1.2)
        sh = RBOX(s, ML, y, CW, h, color, radius=0.04)
        if grad:
            gradient_shape(sh, color, darken(color, 0.15), angle=0.0)
        T(s, ML + Inches(0.3), y + Inches(0.1), CW - Inches(0.6), Inches(0.4),
          text, sz=SZ["subtitle"], c=C["white"], b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["semibold"])
        T(s, ML + Inches(0.3), y + Inches(0.6), CW - Inches(0.6), Inches(0.45),
          sub, sz=SZ["body_sm"], c=C["white"], al=PP_ALIGN.CENTER,
          fn=FONT_W["light"])
    else:
        h = Inches(0.8)
        sh = RBOX(s, ML, y, CW, h, color, text, sz=SZ["subtitle"],
                  tc=C["white"], b=True, radius=0.04)
        if grad:
            gradient_shape(sh, color, darken(color, 0.15), angle=0.0)


def KPIS(s, items, y=None, h=None, shadow=True):
    """KPI 카드 그리드 (v3.6 — 라운드 + 그림자 + 타이포 강화)

    Args:
        items: [{"value": "+30%", "label": "팔로워 성장", "basis": "산출근거"}, ...]
        y: 시작 Y
        h: 카드 높이
        shadow: 그림자 적용 여부
    """
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.8)
    n = len(items)
    gap = 0.15
    total = float(CW / 914400)
    card_w = (total - gap * (n - 1)) / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    h_in = float(h / 914400)
    for i, item in enumerate(items):
        x = ML + Inches((card_w + gap) * i)
        clr = palette[i % len(palette)]
        # 카드 배경 (라운드 + 그림자)
        card_sh = RBOX(s, x, y, Inches(card_w), h + Pt(4),
                       C["card_bg"], radius=0.06)
        card_sh.line.color.rgb = C["card_border"]
        card_sh.line.width = Pt(0.5)
        if shadow:
            add_shadow(card_sh, preset="card")
        # 상단 컬러 바 (카드 상단에 겹침)
        R(s, x, y, Inches(card_w), Pt(4), f=clr)
        # 값 (비율 기반 배치)
        val_y = 0.18
        val_h = h_in * 0.35
        T(s, x + Inches(0.15), y + Inches(val_y), Inches(card_w - 0.3), Inches(val_h),
          item.get("value", ""), sz=28, c=clr, b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        # 라벨
        lbl_y = val_y + val_h
        lbl_h = h_in * 0.2
        T(s, x + Inches(0.15), y + Inches(lbl_y), Inches(card_w - 0.3), Inches(lbl_h),
          item.get("label", ""), sz=SZ["body_sm"], c=C["dark"], b=True,
          al=PP_ALIGN.CENTER, fn=FONT_W["semibold"])
        # 산출근거
        basis = item.get("basis", "")
        if basis:
            basis_y = lbl_y + lbl_h + 0.05
            basis_h = max(h_in - basis_y - 0.05, 0.3)
            T(s, x + Inches(0.1), y + Inches(basis_y), Inches(card_w - 0.2), Inches(basis_h),
              basis, sz=SZ["source"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.3)


def COMPARE(s, left_title, left_items, right_title, right_items,
            y=None, left_color=None, right_color=None):
    """좌우 비교 레이아웃 (AS-IS / TO-BE)

    Args:
        left_title, right_title: 좌/우 제목
        left_items, right_items: 좌/우 항목 리스트
        y: 시작 Y
    """
    if y is None:
        y = Inches(1.2)
    if left_color is None:
        left_color = C["gray"]
    if right_color is None:
        right_color = C["primary"]
    total = float(CW / 914400)
    half = (total - 0.6) / 2   # 0.6" for arrow
    # 좌측
    BOX(s, ML, y, Inches(half), Inches(0.5), left_color,
        left_title, sz=SZ["body"], tc=C["white"], b=True)
    R(s, ML, y + Inches(0.5), Inches(half), Inches(3.0),
      f=C["light"], lc=C["lgray"])
    MT(s, ML + Inches(0.15), y + Inches(0.6),
       Inches(half - 0.3), Inches(2.8),
       left_items, sz=SZ["body_sm"], bul=True)
    # 화살표
    arrow_x = ML + Inches(half)
    T(s, arrow_x, y + Inches(0.5), Inches(0.6), Inches(3.0),
      "→", sz=36, c=C["secondary"], b=True, al=PP_ALIGN.CENTER)
    # 우측
    rx = ML + Inches(half + 0.6)
    BOX(s, rx, y, Inches(half), Inches(0.5), right_color,
        right_title, sz=SZ["body"], tc=C["white"], b=True)
    R(s, rx, y + Inches(0.5), Inches(half), Inches(3.0),
      f=C["light"], lc=C["lgray"])
    MT(s, rx + Inches(0.15), y + Inches(0.6),
       Inches(half - 0.3), Inches(2.8),
       right_items, sz=SZ["body_sm"], bul=True)


def TIMELINE(s, items, y=None, h=None):
    """타임라인 (가로 배치)

    Args:
        items: [("기간", "내용"), ...] 3~6개
        y: 시작 Y
        h: 높이
    """
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.0)
    n = len(items)
    total = float(CW / 914400)
    cell_w = total / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"], C["gold"]]
    # 가로 바
    R(s, ML, y + Inches(0.5), CW, Pt(3), f=C["lgray"])
    for i, (period, content) in enumerate(items):
        x = ML + Inches(cell_w * i)
        clr = palette[i % len(palette)]
        # 마커
        BOX(s, x + Inches(cell_w / 2 - 0.12), y + Inches(0.38),
            Inches(0.24), Inches(0.24), clr, "", sz=8)
        # 기간
        T(s, x, y, Inches(cell_w), Inches(0.35),
          period, sz=SZ["body_sm"], c=clr, b=True, al=PP_ALIGN.CENTER)
        # 내용
        T(s, x + Inches(0.05), y + Inches(0.7), Inches(cell_w - 0.1), h,
          content, sz=SZ["caption"], c=C["dark"], al=PP_ALIGN.CENTER, ls=1.5)


# ═══════════════════════════════════════════════════════════════
#  7. 슬라이드 템플릿
# ═══════════════════════════════════════════════════════════════

def slide_cover(prs, project_name, client_name, year="2026",
                tagline="", company_name="[수행사명]"):
    """표지 슬라이드 (v3.6 — 그라디언트 배경 + 시각 폴리시)"""
    s = new_slide(prs)
    # 그라디언트 배경 (좌하→우상 대각선 다크→블루)
    c1, c2 = GRAD["cover"]()
    gradient_bg(s, c1, c2, angle=225.0)
    # 상단 악센트 라인
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])
    # 좌측 세로 악센트 (미세한 시각 앵커)
    R(s, Inches(0.35), Inches(1.5), Pt(3), Inches(4.0), f=C["secondary"])
    # 프로젝트명
    title_tb = T(s, ML, Inches(1.8), CW, Inches(2.2),
                 project_name, sz=SZ["hero"], c=C["white"], b=True,
                 al=PP_ALIGN.LEFT, fn=FONT_W["bold"])
    set_char_spacing(title_tb, 100)
    # 구분선
    R(s, ML, Inches(4.2), Inches(3), Pt(2), f=C["secondary"])
    # 부제
    if tagline:
        T(s, ML, Inches(4.5), CW, Inches(0.5),
          tagline, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
    # 연도 + 발주처
    T(s, ML, Inches(5.2), CW, Inches(0.4),
      f"{year}  |  {client_name}", sz=SZ["body"], c=C["gray"],
      fn=FONT_W["medium"])
    # 수행사명
    T(s, ML, SH - Inches(0.8), CW, Inches(0.4),
      company_name, sz=SZ["body_sm"], c=C["lgray"],
      fn=FONT_W["light"])
    # 하단 악센트 라인
    R(s, Inches(0), SH - Pt(3), SW, Pt(3), f=C["secondary"])
    return s


def slide_section_divider(prs, num, title, subtitle="", story="",
                          win_theme_key=None, win_themes=None):
    """섹션 구분자 슬라이드 (v3.6 — 그라디언트 + 대형 숫자 아웃라인)"""
    s = new_slide(prs)
    c1, c2 = GRAD["section"]()
    gradient_bg(s, c1, c2, angle=270.0)
    # 상단 악센트 라인
    R(s, Inches(0), Inches(0), SW, Pt(3), f=C["secondary"])
    # 대형 번호 (아웃라인 스타일 — 연한 색)
    T(s, ML, Inches(1.2), Inches(3), Inches(2.2),
      num, sz=110, c=darken(C["secondary"], 0.3), b=True,
      fn=FONT_W["black"])
    # 제목
    title_tb = T(s, ML, Inches(3.5), CW, Inches(0.7),
                 title, sz=SZ["divider"], c=C["white"], b=True,
                 fn=FONT_W["bold"])
    set_char_spacing(title_tb, 80)
    # 부제
    if subtitle:
        T(s, ML, Inches(4.2), CW, Inches(0.4),
          subtitle, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
    # 스토리
    if story:
        T(s, ML, Inches(4.8), CW, Inches(0.4),
          story, sz=SZ["body"], c=C["secondary"])
    # 하단 라인
    R(s, ML, Inches(5.4), Inches(2), Pt(2), f=C["secondary"])
    # Win Theme
    if win_theme_key and win_themes:
        WB(s, win_theme_key, win_themes, ML, Inches(6.0))
    return s


def slide_toc(prs, title, items, pg=None):
    """목차 슬라이드

    Args:
        items: [("01", "HOOK", "설명"), ...]
    """
    s = new_slide(prs)
    TB(s, title, pg)
    y_start = 1.1
    row_h = min(0.55, 5.5 / max(len(items), 1))
    text_h = min(0.3, row_h - 0.08)
    text_pad = (row_h - 0.04 - text_h) / 2
    for i, (num, name, desc) in enumerate(items):
        y = Inches(y_start + row_h * i)
        bgc = C["light"] if i % 2 == 0 else C["white"]
        R(s, ML, y, CW, Inches(row_h - 0.04), f=bgc)
        T(s, ML + Inches(0.2), y + Inches(text_pad), Inches(0.5), Inches(text_h),
          num, sz=SZ["body"], c=C["secondary"], b=True)
        T(s, ML + Inches(0.9), y + Inches(text_pad), Inches(3), Inches(text_h),
          name, sz=SZ["body"], c=C["primary"], b=True)
        T(s, ML + Inches(4.2), y + Inches(text_pad), Inches(7), Inches(text_h),
          desc, sz=SZ["body_sm"], c=C["gray"])
    return s


def slide_exec_summary(prs, title, one_liner, win_themes_dict, kpis, why_us_points):
    """Executive Summary (v3.6 — 그라디언트 하이라이트 + 카드 그림자)"""
    s = new_slide(prs)
    TB(s, title)
    # One Sentence Pitch (그라디언트)
    HIGHLIGHT(s, one_liner, y=Inches(1.1), grad=True)
    # Win Theme (라운드 박스)
    themes = list(win_themes_dict.items())
    colors = [C["primary"], C["secondary"], C["teal"]]
    for i, (key, desc) in enumerate(themes[:3]):
        x = ML + Inches(3.95 * i)
        sh = RBOX(s, x, Inches(2.1), Inches(3.75), Inches(0.55), colors[i % 3],
                  f"Win Theme {i+1}: {desc}", sz=SZ["body_sm"], tc=C["white"],
                  b=True, radius=0.06)
        add_shadow(sh, preset="subtle")
    # KPI
    T(s, ML, Inches(2.85), Inches(3), Inches(0.25),
      "핵심 KPI", sz=SZ["body_sm"], c=C["primary"], b=True,
      fn=FONT_W["semibold"])
    KPIS(s, kpis, y=Inches(3.1), h=Inches(1.5))
    # Why Us
    T(s, ML, Inches(4.8), Inches(3), Inches(0.25),
      "Why Us", sz=SZ["body_sm"], c=C["primary"], b=True,
      fn=FONT_W["semibold"])
    R(s, ML, Inches(5.05), CW, Inches(1.2), f=C["card_bg"], lc=C["card_border"])
    for i, pt in enumerate(why_us_points[:3]):
        x = ML + Inches(3.95 * i) + Inches(0.15)
        T(s, x, Inches(5.15), Inches(3.6), Inches(0.9),
          f"— {pt}", sz=SZ["body_sm"], c=C["dark"], ls=1.5)
    return s


def slide_next_step(prs, headline, steps, contact=""):
    """Next Step / CTA 슬라이드 (v3.6 — 그라디언트 배경)"""
    s = new_slide(prs)
    c1, c2 = GRAD["section"]()
    gradient_bg(s, c1, c2, angle=270.0)
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])
    # NEXT STEP
    ns_tb = T(s, ML, Inches(0.8), CW, Inches(0.5),
              "NEXT STEP", sz=32, c=C["white"], b=True)
    set_char_spacing(ns_tb, 200)
    R(s, ML, Inches(1.3), Inches(1.5), Pt(2), f=C["secondary"])
    # 헤드라인
    T(s, ML, Inches(1.6), CW, Inches(0.4),
      headline, sz=SZ["subtitle"], c=C["lgray"])
    # 스텝 카드
    n = len(steps)
    gap = 0.3
    total = float(CW / 914400)
    card_w = (total - gap * (n - 1)) / n
    card_y = 2.3
    card_h = 3.2
    for i, (step_label, title, desc, clr) in enumerate(steps):
        x = ML + Inches((card_w + gap) * i)
        # 라운드 카드 + 그림자 (v3.6)
        card_sh = RBOX(s, x, Inches(card_y), Inches(card_w), Inches(card_h),
                       clr, radius=0.06)
        add_shadow(card_sh, preset="elevated")
        # 스텝 라벨만 별도 (상단)
        T(s, x, Inches(card_y + 0.2), Inches(card_w), Inches(0.3),
          step_label, sz=SZ["caption"], c=C["white"], al=PP_ALIGN.CENTER,
          fn=FONT_W["medium"])
        T(s, x, Inches(card_y + 0.65), Inches(card_w), Inches(0.5),
          title, sz=20, c=C["white"], b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        T(s, x + Inches(0.2), Inches(card_y + 1.3), Inches(card_w - 0.4), Inches(1.2),
          desc, sz=SZ["body_sm"], c=C["white"], al=PP_ALIGN.CENTER, ls=1.5)
        if i < n - 1:
            T(s, x + Inches(card_w), Inches(card_y + 1.0), Inches(gap), Inches(0.5),
              "→", sz=22, c=C["lgray"], b=True, al=PP_ALIGN.CENTER)
    # 연락처
    if contact:
        T(s, ML, SH - Inches(0.8), CW, Inches(0.5),
          f"Contact: {contact}", sz=SZ["body_sm"], c=C["lgray"], al=PP_ALIGN.CENTER)
    return s


def slide_closing(prs, message="감사합니다", tagline="",
                  project_title="", contact=""):
    """마지막 감사 슬라이드 (v3.6 — 그라디언트 + 시각 폴리시)"""
    s = new_slide(prs)
    c1, c2 = GRAD["closing"]()
    gradient_bg(s, c1, c2, angle=225.0)
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])
    # 좌측 세로 악센트
    R(s, Inches(0.35), Inches(2.2), Pt(3), Inches(2.0), f=C["secondary"])
    # 메인 메시지
    msg_tb = T(s, ML, Inches(2.5), CW, Inches(1.0),
               message, sz=SZ["hero"], c=C["white"], b=True, al=PP_ALIGN.LEFT,
               fn=FONT_W["bold"])
    set_char_spacing(msg_tb, 100)
    # 구분선
    R(s, ML, Inches(3.8), Inches(2), Pt(2), f=C["secondary"])
    # 태그라인
    if tagline:
        T(s, ML, Inches(4.1), CW, Inches(0.4),
          tagline, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
    # 프로젝트
    if project_title:
        T(s, ML, Inches(4.7), CW, Inches(0.4),
          project_title, sz=SZ["body"], c=C["gray"],
          fn=FONT_W["medium"])
    # 연락처
    if contact:
        T(s, ML, SH - Inches(0.8), CW, Inches(0.4),
          contact, sz=SZ["body_sm"], c=C["lgray"])
    # 하단 라인
    R(s, Inches(0), SH - Pt(3), SW, Pt(3), f=C["secondary"])
    return s


# ═══════════════════════════════════════════════════════════════
#  8. 유틸리티
# ═══════════════════════════════════════════════════════════════

def save_pptx(prs, output_path):
    """프레젠테이션 저장"""
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    n = len(prs.slides)
    print(f"생성 완료: {output_path} ({n}장)")
    return output_path


# ═══════════════════════════════════════════════════════════════
#  9. 레이아웃 Zone 시스템 — 겹침 방지 표준 영역
# ═══════════════════════════════════════════════════════════════

# 표준 영역 (TB 타이틀바 포함 기준)
Z = {
    "tb_y":  0,        # 타이틀바 시작
    "tb_h":  0.88,     # 타이틀바 높이 (TB 함수 기준)
    "ct_y":  1.1,      # 콘텐츠 시작
    "ct_h":  5.4,      # 콘텐츠 높이
    "ct_b":  6.5,      # 콘텐츠 하단
    "ft_y":  6.7,      # 푸터 시작
    "ft_h":  0.8,      # 푸터 높이
}

# 안전 간격
GAP = 0.2    # 요소 간 수직 간격
CGAP = 0.15  # 컬럼 간 수평 간격
CW_IN = float(CW / 914400)   # CW in inches
ML_IN = float(ML / 914400)   # ML in inches


def _cols(n, gap=CGAP):
    """N등분 컬럼 너비 계산 (inches)"""
    return (CW_IN - gap * (n - 1)) / n


# ═══════════════════════════════════════════════════════════════
#  10. LAYOUTS 데이터 — 20가지 레이아웃 프리셋
# ═══════════════════════════════════════════════════════════════

"""
레이아웃 데이터 규격:
- 각 레이아웃은 name, desc, zones 를 가짐
- zones: [{"id": str, "x": float, "y": float, "w": float, "h": float, "role": str}]
  - x, y, w, h 는 모두 인치(inches) 단위
  - role: "header"|"body"|"image"|"table"|"card"|"kpi"|"footer"
- 모든 위치는 TB() 이후 콘텐츠 영역 기준으로 사전 계산됨
"""

LAYOUTS = {

    # ── 1. 풀바디 ────────────────────────────────────────────────
    "FULL_BODY": {
        "desc": "타이틀 + 전체 너비 본문 텍스트",
        "zones": [
            {"id": "body", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 2. 하이라이트 + 본문 ──────────────────────────────────────
    "HIGHLIGHT_BODY": {
        "desc": "강조 메시지 + 본문",
        "zones": [
            {"id": "highlight", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": 0.8, "role": "header"},
            {"id": "body", "x": ML_IN, "y": Z["ct_y"] + 1.0, "w": CW_IN, "h": Z["ct_h"] - 1.0, "role": "body"},
        ],
    },

    # ── 3. 2단 컬럼 ────────────────────────────────────────────────
    "TWO_COL": {
        "desc": "타이틀 + 좌우 2단 레이아웃",
        "zones": [
            {"id": "left",  "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(2), "h": Z["ct_h"], "role": "body"},
            {"id": "right", "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"],
             "w": _cols(2), "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 4. 3단 컬럼 ────────────────────────────────────────────────
    "THREE_COL": {
        "desc": "타이틀 + 3단 비교 레이아웃",
        "zones": [
            {"id": "col1", "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
            {"id": "col2", "x": ML_IN + (_cols(3) + CGAP), "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
            {"id": "col3", "x": ML_IN + (_cols(3) + CGAP) * 2, "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 5. 4단 컬럼 ────────────────────────────────────────────────
    "FOUR_COL": {
        "desc": "타이틀 + 4단 카드 레이아웃",
        "zones": [
            {"id": f"col{i+1}",
             "x": ML_IN + (_cols(4) + CGAP) * i, "y": Z["ct_y"],
             "w": _cols(4), "h": Z["ct_h"], "role": "card"}
            for i in range(4)
        ],
    },

    # ── 6. 좌우 비교 (AS-IS / TO-BE) ─────────────────────────────
    "COMPARE_LR": {
        "desc": "좌우 비교 (Before/After)",
        "zones": [
            {"id": "left_header",  "x": ML_IN, "y": Z["ct_y"],
             "w": (CW_IN - 0.6) / 2, "h": 0.5, "role": "header"},
            {"id": "left_body",    "x": ML_IN, "y": Z["ct_y"] + 0.5,
             "w": (CW_IN - 0.6) / 2, "h": Z["ct_h"] - 0.5, "role": "body"},
            {"id": "arrow",        "x": ML_IN + (CW_IN - 0.6) / 2, "y": Z["ct_y"],
             "w": 0.6, "h": Z["ct_h"], "role": "body"},
            {"id": "right_header", "x": ML_IN + (CW_IN - 0.6) / 2 + 0.6, "y": Z["ct_y"],
             "w": (CW_IN - 0.6) / 2, "h": 0.5, "role": "header"},
            {"id": "right_body",   "x": ML_IN + (CW_IN - 0.6) / 2 + 0.6, "y": Z["ct_y"] + 0.5,
             "w": (CW_IN - 0.6) / 2, "h": Z["ct_h"] - 0.5, "role": "body"},
        ],
    },

    # ── 7. 하이라이트 + 3단 카드 ─────────────────────────────────
    "HIGHLIGHT_THREE_CARD": {
        "desc": "강조 메시지 + 3단 카드",
        "zones": [
            {"id": "highlight", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": 0.8, "role": "header"},
            {"id": "card1", "x": ML_IN, "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
            {"id": "card2", "x": ML_IN + (_cols(3) + CGAP), "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
            {"id": "card3", "x": ML_IN + (_cols(3) + CGAP) * 2, "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
        ],
    },

    # ── 8. KPI 카드 그리드 ──────────────────────────────────────
    "KPI_GRID": {
        "desc": "타이틀 + KPI 카드 + 산출근거",
        "zones": [
            {"id": "kpi_row", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 2.0, "role": "kpi"},
            {"id": "detail", "x": ML_IN, "y": Z["ct_y"] + 2.3,
             "w": CW_IN, "h": Z["ct_h"] - 2.3, "role": "body"},
        ],
    },

    # ── 9. 프로세스 플로우 + 설명 ────────────────────────────────
    "PROCESS_DESC": {
        "desc": "프로세스 플로우 + 하단 상세 설명",
        "zones": [
            {"id": "flow", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 1.2, "role": "header"},
            {"id": "flow_desc", "x": ML_IN, "y": Z["ct_y"] + 1.4,
             "w": CW_IN, "h": 0.8, "role": "body"},
            {"id": "detail", "x": ML_IN, "y": Z["ct_y"] + 2.4,
             "w": CW_IN, "h": Z["ct_h"] - 2.4, "role": "body"},
        ],
    },

    # ── 10. 타임라인 + 하단 설명 ─────────────────────────────────
    "TIMELINE_DESC": {
        "desc": "타임라인 + 하단 본문",
        "zones": [
            {"id": "timeline", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 2.2, "role": "header"},
            {"id": "body", "x": ML_IN, "y": Z["ct_y"] + 2.5,
             "w": CW_IN, "h": Z["ct_h"] - 2.5, "role": "body"},
        ],
    },

    # ── 11. 피라미드 + 우측 설명 ─────────────────────────────────
    "PYRAMID_DESC": {
        "desc": "좌측 피라미드 + 우측 설명",
        "zones": [
            {"id": "pyramid", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN * 0.45, "h": Z["ct_h"], "role": "body"},
            {"id": "desc", "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"],
             "w": CW_IN * 0.5, "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 12. 2×2 매트릭스 + 하단 설명 ─────────────────────────────
    "MATRIX_DESC": {
        "desc": "매트릭스 + 하단 시사점",
        "zones": [
            {"id": "matrix", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 4.0, "role": "body"},
            {"id": "insight", "x": ML_IN, "y": Z["ct_y"] + 4.2,
             "w": CW_IN, "h": Z["ct_h"] - 4.2, "role": "body"},
        ],
    },

    # ── 13. 이미지 갤러리 (3×2 그리드) ─────────────────────────────
    "GALLERY_3x2": {
        "desc": "3열 2행 이미지 갤러리 + 캡션",
        "zones": [
            {"id": f"img_{r}_{c}",
             "x": ML_IN + (_cols(3) + CGAP) * c,
             "y": Z["ct_y"] + (2.5 + GAP) * r,
             "w": _cols(3), "h": 2.3, "role": "image"}
            for r in range(2) for c in range(3)
        ],
    },

    # ── 14. 키비주얼 (좌측 이미지 + 우측 텍스트) ──────────────────
    "KEY_VISUAL": {
        "desc": "좌측 대형 이미지 + 우측 텍스트",
        "zones": [
            {"id": "image", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN * 0.45, "h": Z["ct_h"], "role": "image"},
            {"id": "title", "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"],
             "w": CW_IN * 0.5, "h": 0.6, "role": "header"},
            {"id": "body",  "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"] + 0.8,
             "w": CW_IN * 0.5, "h": Z["ct_h"] - 0.8, "role": "body"},
        ],
    },

    # ── 15. 테이블 + 인사이트 ──────────────────────────────────────
    "TABLE_INSIGHT": {
        "desc": "데이터 테이블 + 하단 인사이트 박스",
        "zones": [
            {"id": "table", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 3.5, "role": "table"},
            {"id": "insight", "x": ML_IN, "y": Z["ct_y"] + 3.8,
             "w": CW_IN, "h": Z["ct_h"] - 3.8, "role": "body"},
        ],
    },

    # ── 16. 프로그램 카드 (이미지+내용+포인트) ─────────────────────
    "PROGRAM_CARD_3": {
        "desc": "3단 프로그램 카드 (이미지 + 본문 + 포인트)",
        "zones": [
            {"id": f"card{i+1}_img",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"],
             "w": _cols(3), "h": 2.0, "role": "image"}
            for i in range(3)
        ] + [
            {"id": f"card{i+1}_body",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"] + 2.1,
             "w": _cols(3), "h": 2.3, "role": "body"}
            for i in range(3)
        ] + [
            {"id": f"card{i+1}_point",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"] + 4.5,
             "w": _cols(3), "h": 0.9, "role": "footer"}
            for i in range(3)
        ],
    },

    # ── 17. 4분할 공간 ──────────────────────────────────────────
    "QUAD_GRID": {
        "desc": "2×2 이미지 그리드 + 각 캡션",
        "zones": [
            {"id": f"quad_{r}_{c}",
             "x": ML_IN + (_cols(2) + CGAP) * c,
             "y": Z["ct_y"] + (2.6 + GAP) * r,
             "w": _cols(2), "h": 2.4, "role": "image"}
            for r in range(2) for c in range(2)
        ],
    },

    # ── 18. 조직도 (3단 계층) ──────────────────────────────────────
    "ORG_CHART": {
        "desc": "PM + 감독 + 팀 3단 계층",
        "zones": [
            {"id": "pm", "x": ML_IN + CW_IN * 0.35, "y": Z["ct_y"],
             "w": CW_IN * 0.3, "h": 1.2, "role": "card"},
            {"id": "dir1", "x": ML_IN, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir2", "x": ML_IN + (_cols(4) + CGAP), "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir3", "x": ML_IN + (_cols(4) + CGAP) * 2, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir4", "x": ML_IN + (_cols(4) + CGAP) * 3, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "team_row", "x": ML_IN, "y": Z["ct_y"] + 3.2,
             "w": CW_IN, "h": Z["ct_h"] - 3.2, "role": "body"},
        ],
    },

    # ── 19. 리스크 카드 (2열 × 3단 대응) ──────────────────────────
    "RISK_CARD": {
        "desc": "좌우 리스크 + 3단 대응 방안",
        "zones": [
            {"id": "risk1_title", "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(2), "h": 0.5, "role": "header"},
            {"id": "risk1_body",  "x": ML_IN, "y": Z["ct_y"] + 0.6,
             "w": _cols(2), "h": 2.0, "role": "body"},
            {"id": "risk1_resp",  "x": ML_IN, "y": Z["ct_y"] + 2.8,
             "w": _cols(2), "h": Z["ct_h"] - 2.8, "role": "body"},
            {"id": "risk2_title", "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"],
             "w": _cols(2), "h": 0.5, "role": "header"},
            {"id": "risk2_body",  "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"] + 0.6,
             "w": _cols(2), "h": 2.0, "role": "body"},
            {"id": "risk2_resp",  "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"] + 2.8,
             "w": _cols(2), "h": Z["ct_h"] - 2.8, "role": "body"},
        ],
    },

    # ── 20. 연간 간트 차트 (12개월) ───────────────────────────────
    "GANTT": {
        "desc": "월별 간트 차트 (좌측 카테고리 + 12개월 그리드)",
        "zones": [
            {"id": "categories", "x": ML_IN, "y": Z["ct_y"],
             "w": 2.0, "h": Z["ct_h"], "role": "body"},
            {"id": "grid", "x": ML_IN + 2.1, "y": Z["ct_y"],
             "w": CW_IN - 2.1, "h": Z["ct_h"], "role": "table"},
        ],
    },
}


def get_zones(layout_name):
    """레이아웃 프리셋의 zone 목록 반환 (인치 단위 dict 리스트)

    Usage:
        zones = get_zones("TWO_COL")
        left = zones["left"]   # {"x": ..., "y": ..., "w": ..., "h": ..., "role": ...}
        right = zones["right"]
    Returns:
        dict[str, dict] — id를 key로 하는 zone 딕셔너리
    """
    layout = LAYOUTS.get(layout_name)
    if not layout:
        raise ValueError(f"Unknown layout: {layout_name}. "
                         f"Available: {list(LAYOUTS.keys())}")
    return {z["id"]: z for z in layout["zones"]}


def zone_to_inches(z):
    """zone dict → (Inches(x), Inches(y), Inches(w), Inches(h)) 튜플"""
    return Inches(z["x"]), Inches(z["y"]), Inches(z["w"]), Inches(z["h"])


# ═══════════════════════════════════════════════════════════════
#  11. 추가 도식화 헬퍼 — 레퍼런스 분석 기반
# ═══════════════════════════════════════════════════════════════

def GRID(s, items, cols=3, y=None, h=None, gap=CGAP, shadow=True):
    """N×M 카드 그리드 (v3.6 — 라운드 헤더 + 그림자)

    Args:
        items: [{"title": "제목", "body": "본문" or ["줄1","줄2"], "color": RGBColor}, ...]
        cols: 열 수 (2~4)
        y: 시작 Y (기본 1.1")
        h: 카드 높이 (기본 자동)
        gap: 간격
        shadow: 그림자 적용 여부
    """
    if y is None:
        y = Inches(Z["ct_y"])
    n = len(items)
    rows = (n + cols - 1) // cols
    col_w = _cols(cols, gap)
    if h is None:
        card_h = min(2.5, (Z["ct_h"] - gap * (rows - 1)) / rows)
    else:
        card_h = float(h / 914400)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"],
               C["green"], C["gold"]]
    for idx, item in enumerate(items):
        c_idx = idx % cols
        r_idx = idx // cols
        x = ML + Inches((col_w + gap) * c_idx)
        iy = y + Inches((card_h + gap) * r_idx)
        clr = item.get("color", palette[idx % len(palette)])
        title = item.get("title", "")
        body = item.get("body", "")
        # 카드 배경 (그림자)
        card_bg_sh = R(s, x, iy, Inches(col_w), Inches(card_h),
                       f=C["card_bg"], lc=C["card_border"], lw=0.5)
        if shadow:
            add_shadow(card_bg_sh, preset="card")
        # 헤더
        RBOX(s, x, iy, Inches(col_w), Inches(0.45), clr,
             title, sz=SZ["body"], tc=C["white"], b=True, radius=0.0)
        # 바디
        body_y = iy + Inches(0.45)
        body_h = Inches(card_h - 0.45)
        if isinstance(body, list):
            MT(s, x + Inches(0.1), body_y + Inches(0.05),
               Inches(col_w - 0.2), body_h - Inches(0.1),
               body, sz=SZ["body_sm"], bul=True)
        else:
            T(s, x + Inches(0.1), body_y + Inches(0.05),
              Inches(col_w - 0.2), body_h - Inches(0.1),
              body, sz=SZ["body_sm"])


def STAT_ROW(s, items, y=None, h=None, shadow=True):
    """통계/수치 강조 행 (v3.6 — 그림자 + 타이포 강화)

    Args:
        items: [{"value": "87%", "label": "달성률", "color": RGBColor}, ...]
        y: 시작 Y
        h: 높이
        shadow: 그림자 적용 여부
    """
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.2)
    n = len(items)
    col_w = _cols(n)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    for i, item in enumerate(items):
        x = ML + Inches((col_w + CGAP) * i)
        clr = item.get("color", palette[i % len(palette)])
        # 카드 배경 (그림자)
        card_sh = R(s, x, y, Inches(col_w), h + Pt(4),
                    f=C["card_bg"], lc=C["card_border"], lw=0.5)
        if shadow:
            add_shadow(card_sh, preset="subtle")
        # 상단 라인
        R(s, x, y, Inches(col_w), Pt(4), f=clr)
        # 수치
        T(s, x, y + Inches(0.12), Inches(col_w), Inches(0.6),
          item.get("value", ""), sz=32, c=clr, b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        # 라벨
        T(s, x, y + Inches(0.72), Inches(col_w), Inches(0.35),
          item.get("label", ""), sz=SZ["body_sm"], c=C["dark"], al=PP_ALIGN.CENTER,
          fn=FONT_W["semibold"])


def GANTT_CHART(s, categories, months, data, y=None, colors=None):
    """간트 차트 (좌측 카테고리 + 월별 컬러 바)

    Args:
        categories: ["기획", "실행", "보고"]
        months: ["3월", "4월", "5월", ...]
        data: [[1,1,0,0,...], [0,1,1,1,...], ...] — 각 카테고리별 활성 월 (1/0)
        y: 시작 Y
        colors: 카테고리별 색상
    """
    if y is None:
        y = Inches(Z["ct_y"])
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"],
                   C["green"], C["gold"]]
        colors = [palette[i % len(palette)] for i in range(len(categories))]
    cat_w = 2.0
    n_months = len(months)
    month_w = (CW_IN - cat_w - 0.1) / n_months
    row_h = min(0.5, (Z["ct_h"] - 0.5) / max(len(categories), 1))
    # 헤더
    BOX(s, ML, y, Inches(cat_w), Inches(0.4), C["primary"],
        "구분", sz=SZ["body_sm"], tc=C["white"], b=True)
    for j, m in enumerate(months):
        mx = ML + Inches(cat_w + 0.1 + month_w * j)
        BOX(s, mx, y, Inches(month_w - 0.02), Inches(0.4), C["primary"],
            m, sz=SZ["caption"], tc=C["white"], b=True)
    # 데이터 행
    for i, cat in enumerate(categories):
        ry = y + Inches(0.45 + row_h * i)
        bgc = C["light"] if i % 2 == 0 else C["white"]
        # 카테고리명
        R(s, ML, ry, Inches(cat_w), Inches(row_h - 0.02), f=bgc)
        T(s, ML + Inches(0.1), ry + Inches(0.02), Inches(cat_w - 0.2),
          Inches(row_h - 0.06), cat, sz=SZ["body_sm"], c=C["dark"], b=True)
        # 월별 바
        for j in range(n_months):
            mx = ML + Inches(cat_w + 0.1 + month_w * j)
            if i < len(data) and j < len(data[i]) and data[i][j]:
                R(s, mx, ry + Inches(0.06), Inches(month_w - 0.02),
                  Inches(row_h - 0.14), f=colors[i])
            else:
                R(s, mx, ry, Inches(month_w - 0.02),
                  Inches(row_h - 0.02), f=bgc)


def ORG(s, pm, directors, teams=None, y=None):
    """조직도 (PM + 감독 + 팀원)

    Args:
        pm: {"name": "PM명", "role": "프로젝트 매니저", "detail": "상세"}
        directors: [{"name": "감독1", "role": "역할"}, ...]
        teams: [{"name": "팀원1", "role": "역할"}, ...] (선택)
        y: 시작 Y
    """
    if y is None:
        y = Inches(Z["ct_y"])
    # PM 박스
    pm_w = 3.0
    pm_x = ML_IN + (CW_IN - pm_w) / 2
    pm_h = 1.2 if pm.get("detail") else 1.0
    BOX(s, Inches(pm_x), y, Inches(pm_w), Inches(pm_h), C["primary"],
        f"{pm.get('name', 'PM')}\n{pm.get('role', '')}", sz=SZ["body"], tc=C["white"], b=True)
    if pm.get("detail"):
        T(s, Inches(pm_x), y + Inches(0.88), Inches(pm_w), Inches(0.25),
          pm["detail"], sz=SZ["source"], c=C["white"], al=PP_ALIGN.CENTER)
    # 연결선
    line_y = y + Inches(pm_h)
    R(s, Inches(pm_x + pm_w / 2 - 0.01), line_y, Pt(2), Inches(0.3), f=C["lgray"])
    # Directors
    dir_y = line_y + Inches(0.3)
    n = len(directors)
    dir_w = _cols(n, CGAP)
    R(s, ML, dir_y, CW, Pt(2), f=C["lgray"])
    for i, d in enumerate(directors):
        dx = ML + Inches((dir_w + CGAP) * i)
        BOX(s, dx, dir_y + Inches(0.1), Inches(dir_w), Inches(0.9), C["secondary"],
            f"{d.get('name', '')}\n{d.get('role', '')}", sz=SZ["body_sm"], tc=C["white"], b=True)
    # Teams
    if teams:
        team_y = dir_y + Inches(1.3)
        n_t = len(teams)
        t_w = _cols(n_t, 0.1)
        for i, t in enumerate(teams):
            tx = ML + Inches((t_w + 0.1) * i)
            OBOX(s, tx, team_y, Inches(t_w), Inches(0.7),
                 f"{t.get('name', '')} — {t.get('role', '')}",
                 sz=SZ["caption"], tc=C["dark"], lc=C["lgray"])


def ICON_CARDS(s, items, y=None, h=None):
    """아이콘 + 텍스트 카드 행 (포인트 태그)

    Args:
        items: [{"icon": "★", "title": "제목", "desc": "설명"}, ...]
        y: 시작 Y
        h: 카드 높이
    """
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.5)
    n = len(items)
    col_w = _cols(n)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    for i, item in enumerate(items):
        x = ML + Inches((col_w + CGAP) * i)
        clr = palette[i % len(palette)]
        # 아이콘 서클
        BOX(s, x + Inches(col_w / 2 - 0.3), y, Inches(0.6), Inches(0.6), clr,
            item.get("icon", "●"), sz=24, tc=C["white"], b=True)
        # 타이틀
        T(s, x, y + Inches(0.7), Inches(col_w), Inches(0.3),
          item.get("title", ""), sz=SZ["body"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
        # 설명
        T(s, x + Inches(0.1), y + Inches(1.0), Inches(col_w - 0.2),
          h - Inches(1.0),
          item.get("desc", ""), sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  12. 시각화 헬퍼
# ═══════════════════════════════════════════════════════════════

def IMG_PH(s, x, y, w, h, label="이미지 영역"):
    """이미지 플레이스홀더 — 회색 박스 + 아이콘 + 라벨"""
    R(s, x, y, w, h, f=C["light"], lc=C["lgray"])
    h_in = float(h / 914400)
    T(s, x, y + Inches(0.05), w, Inches(h_in * 0.5),
      "[IMG]", sz=28, c=C["lgray"], al=PP_ALIGN.CENTER)
    T(s, x, y + Inches(h_in * 0.6), w, Inches(h_in * 0.3),
      label, sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER)


def PROGRESS_BAR(s, x, y, w, label, value, max_val=100, color=None, show_pct=True):
    """프로그레스 바 — 라벨 + 바 + 수치"""
    if color is None:
        color = C["secondary"]
    bar_h = Inches(0.22)
    # 라벨
    T(s, x, y, Inches(2.5), Inches(0.25),
      label, sz=SZ["body_sm"], c=C["dark"], b=True)
    # 배경 바
    bar_x = x + Inches(2.6)
    bar_w = w - Inches(3.5)
    R(s, bar_x, y + Inches(0.02), bar_w, bar_h, f=C["light"])
    # 채움 바
    fill_w = Inches(float(bar_w / 914400) * min(value / max_val, 1.0))
    if float(fill_w / 914400) > 0.05:
        R(s, bar_x, y + Inches(0.02), fill_w, bar_h, f=color)
    # 수치
    txt = f"{value}%" if show_pct else str(value)
    T(s, x + w - Inches(0.8), y, Inches(0.8), Inches(0.25),
      txt, sz=SZ["body_sm"], c=color, b=True, al=PP_ALIGN.RIGHT)


def METRIC_CARD(s, x, y, w, h, value, label, sub="", color=None, shadow=True):
    """메트릭 카드 (v3.6 — 그림자 + 라운드)"""
    if color is None:
        color = C["primary"]
    h_in = float(h / 914400)
    # 카드 배경 (그림자)
    card_sh = R(s, x, y, w, h, f=C["card_bg"], lc=C["card_border"], lw=0.5)
    if shadow:
        add_shadow(card_sh, preset="subtle")
    # 상단 컬러 바
    R(s, x, y, w, Pt(4), f=color)
    # 비율 기반 배치 — 카드 높이에 비례
    val_sz = max(24, min(40, int(h_in * 24)))
    T(s, x, y + Inches(h_in * 0.08), w, Inches(h_in * 0.38),
      str(value), sz=val_sz, c=color, b=True, al=PP_ALIGN.CENTER)
    T(s, x, y + Inches(h_in * 0.46), w, Inches(h_in * 0.22),
      label, sz=SZ["body_sm"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
    if sub:
        T(s, x + Inches(0.08), y + Inches(h_in * 0.70), w - Inches(0.16),
          Inches(h_in * 0.28),
          sub, sz=SZ["caption"], c=C["gray"], al=PP_ALIGN.CENTER)


def STEP_ARROW(s, items, y=None, h=None):
    """화살표 스텝 다이어그램 — 숫자 원 + 제목 + 설명 (가로)"""
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.8)
    n = len(items)
    total = float(CW / 914400)
    arrow_w = 0.3
    item_w = (total - arrow_w * (n - 1)) / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
    for i, (num, title, desc) in enumerate(items):
        clr = palette[i % len(palette)]
        x = ML + Inches((item_w + arrow_w) * i)
        # 원형 숫자
        circle_sz = 0.5
        cx = float(x / 914400) + item_w / 2 - circle_sz / 2
        BOX(s, Inches(cx), y, Inches(circle_sz), Inches(circle_sz), clr,
            str(num), sz=18, tc=C["white"], b=True)
        # 제목
        T(s, x, y + Inches(0.6), Inches(item_w), Inches(0.35),
          title, sz=SZ["body"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
        # 설명
        T(s, x + Inches(0.1), y + Inches(0.95), Inches(item_w - 0.2), h - Inches(1.0),
          desc, sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.4)
        # 화살표
        if i < n - 1:
            ax = x + Inches(item_w)
            T(s, ax, y + Inches(0.1), Inches(arrow_w), Inches(0.5),
              "→", sz=20, c=C["lgray"], b=True, al=PP_ALIGN.CENTER)


def DONUT_LABEL(s, x, y, w, value, label, color=None):
    """도넛 차트 스타일 라벨 — 원형 + 큰 숫자 + 라벨 (세로 배치)"""
    if color is None:
        color = C["primary"]
    # 원형 배경
    circle_d = min(float(w / 914400) * 0.7, 1.2)
    cx = float(x / 914400) + (float(w / 914400) - circle_d) / 2
    BOX(s, Inches(cx), y, Inches(circle_d), Inches(circle_d), color,
        str(value), sz=24, tc=C["white"], b=True)
    # 라벨
    T(s, x, y + Inches(circle_d + 0.1), w, Inches(0.3),
      label, sz=SZ["body_sm"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  13. 유틸리티
# ═══════════════════════════════════════════════════════════════

def list_layouts():
    """사용 가능한 모든 레이아웃 목록 출력"""
    for name, layout in LAYOUTS.items():
        n_zones = len(layout["zones"])
        print(f"  {name:25s} — {layout['desc']}  ({n_zones} zones)")


# ═══════════════════════════════════════════════════════════════
#  14. v3.5 — VStack 자동 수직 스택
# ═══════════════════════════════════════════════════════════════

class VStack:
    """수직 자동 스택 — Y좌표를 자동 계산하여 겹침/공백 방지

    Usage:
        v = VStack()
        HIGHLIGHT(s, "메시지", y=v.next(0.8))
        COLS(s, items, y=v.next(3.5), h=Inches(3.5))
        MT(s, ML, v.next(1.4), CW, Inches(1.4), lines, bul=True)
    """

    def __init__(self, y_start=None, gap=GAP):
        self.y = y_start if y_start is not None else Z["ct_y"]
        self.gap = gap

    def next(self, height):
        """다음 요소의 Y위치(Inches) 반환 후 커서를 height+gap만큼 이동"""
        y = self.y
        self.y += height + self.gap
        return Inches(y)

    def next_raw(self, height):
        """다음 요소의 Y위치(float, inches) 반환 — Inches 래핑 없이"""
        y = self.y
        self.y += height + self.gap
        return y

    def skip(self, amount=0.2):
        """추가 여백 삽입"""
        self.y += amount
        return self

    def peek(self):
        """현재 Y위치 반환 (커서 이동 없음)"""
        return Inches(self.y)

    def peek_raw(self):
        """현재 Y위치 반환 (float, 커서 이동 없음)"""
        return self.y

    @property
    def remaining(self):
        """남은 콘텐츠 영역 높이 (inches)"""
        return Z["ct_b"] - self.y

    @property
    def is_full(self):
        """남은 공간이 0.5" 미만이면 True"""
        return self.remaining < 0.5


# ═══════════════════════════════════════════════════════════════
#  15. v3.5 — 라운드 코너 박스
# ═══════════════════════════════════════════════════════════════

def RBOX(s, l, t, w, h, f, text="", sz=13, tc=None, b=False, radius=0.12):
    """라운드 코너 텍스트 박스 — 카드/배지에 부드러운 인상

    Args:
        radius: 코너 반경 비율 (0.0~0.5). 0.12 = 적당히 둥글게
    """
    if tc is None:
        tc = C["white"]
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    # 라운드 비율 설정 (0.0 = 직각, 0.5 = 완전 원형)
    sh.adjustments[0] = min(radius, 0.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = f
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def ORBOX(s, l, t, w, h, text="", sz=13, tc=None, b=False, lc=None, radius=0.12):
    """라운드 아웃라인 박스 — 배경 투명, 테두리 + 라운드"""
    if tc is None:
        tc = C["dark"]
    if lc is None:
        lc = C["primary"]
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = min(radius, 0.5)
    sh.fill.background()
    sh.line.color.rgb = lc
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


# ═══════════════════════════════════════════════════════════════
#  16. v3.5 — 미세 그림자 + 반투명 오버레이
# ═══════════════════════════════════════════════════════════════

# 그림자 프리셋 (v3.6)
SHADOW = {
    "subtle":   {"blur_pt": 2, "offset_pt": 1, "alpha": 78000},   # 은은한 떠있는 느낌
    "normal":   {"blur_pt": 3, "offset_pt": 2, "alpha": 65000},   # 기본 깊이감
    "elevated": {"blur_pt": 5, "offset_pt": 3, "alpha": 55000},   # 강한 부유감
    "card":     {"blur_pt": 4, "offset_pt": 2, "alpha": 72000},   # 카드 전용 (부드러운)
}


def add_shadow(shape, blur_pt=3, offset_pt=2, direction=2700000, alpha=60000,
               preset=None):
    """미세 그림자 — 카드/박스에 깊이감 부여

    Args:
        blur_pt: 블러 반경 (pt)
        offset_pt: 그림자 거리 (pt)
        direction: 각도 (2700000 = 우하단, 단위: 60000분의 1도)
        alpha: 불투명도 (0=불투명, 100000=완전투명)
        preset: SHADOW 프리셋 키 ("subtle"/"normal"/"elevated"/"card")
    """
    if preset and preset in SHADOW:
        p = SHADOW[preset]
        blur_pt = p["blur_pt"]
        offset_pt = p["offset_pt"]
        alpha = p["alpha"]
    try:
        from lxml import etree
    except ImportError:
        return shape
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    # 기존 effectLst 제거
    for old in spPr.findall(f'{{{ns}}}effectLst'):
        spPr.remove(old)
    effectLst = etree.SubElement(spPr, f'{{{ns}}}effectLst')
    outerShdw = etree.SubElement(effectLst, f'{{{ns}}}outerShdw',
                                  blurRad=str(blur_pt * 12700),
                                  dist=str(offset_pt * 12700),
                                  dir=str(direction), algn='tl')
    srgb = etree.SubElement(outerShdw, f'{{{ns}}}srgbClr', val='000000')
    etree.SubElement(srgb, f'{{{ns}}}alpha', val=str(alpha))
    return shape


def OVERLAY(s, l, t, w, h, color, alpha=50000):
    """반투명 오버레이 — 이미지 위 텍스트 가독성 확보

    Args:
        color: 오버레이 색상
        alpha: 0=완전불투명, 100000=완전투명 (50000=반투명)
    """
    sh = R(s, l, t, w, h, f=color)
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        # shape의 spPr 내 solidFill/srgbClr에 alpha 추가
        spPr = sh._element.spPr
        srgbClr = spPr.find(f'.//{{{ns}}}srgbClr')
        if srgbClr is not None:
            for old in srgbClr.findall(f'{{{ns}}}alpha'):
                srgbClr.remove(old)
            etree.SubElement(srgbClr, f'{{{ns}}}alpha', val=str(alpha))
    except (ImportError, AttributeError):
        pass
    return sh


# ═══════════════════════════════════════════════════════════════
#  17. v3.5 — 구분선 / 악센트 요소
# ═══════════════════════════════════════════════════════════════

def DIVIDER(s, y, style="line", color=None, w=None):
    """수평 구분선

    Args:
        style: "line" (실선), "thick" (두꺼운 선), "double" (이중선)
        color: 색상 (기본 lgray)
        w: 너비 (기본 CW)
    """
    if color is None:
        color = C["lgray"]
    if w is None:
        w = CW
    if style == "line":
        R(s, ML, Inches(y), w, Pt(1), f=color)
    elif style == "thick":
        R(s, ML, Inches(y), w, Pt(3), f=color)
    elif style == "double":
        R(s, ML, Inches(y), w, Pt(1), f=color)
        R(s, ML, Inches(y + 0.06), w, Pt(1), f=color)


def ACCENT_LINE(s, x, y, h, color=None, w_pt=3):
    """좌측 악센트 라인 — 인용문/강조 블록 좌측 수직선

    Args:
        x, y: 위치 (inches)
        h: 높이 (inches)
        w_pt: 선 두께 (pt)
    """
    if color is None:
        color = C["secondary"]
    R(s, Inches(x), Inches(y), Pt(w_pt), Inches(h), f=color)


# ═══════════════════════════════════════════════════════════════
#  18. v3.5 — 인용문 / 번호 리스트
# ═══════════════════════════════════════════════════════════════

def QUOTE(s, text, author="", y=None, color=None, style="modern"):
    """인용문 블록

    Args:
        text: 인용문 텍스트
        author: 출처/저자
        y: Y 위치 (Inches 객체 또는 None)
        style: "modern" (좌측 악센트 라인) / "box" (박스형)
    """
    if y is None:
        y = Inches(Z["ct_y"])
    if color is None:
        color = C["secondary"]
    if style == "modern":
        total_h = 1.0 if not author else 1.3
        # 좌측 악센트 라인
        R(s, ML, y, Pt(4), Inches(total_h), f=color)
        # 인용문 텍스트
        T(s, ML + Inches(0.3), y + Inches(0.08), CW - Inches(0.3), Inches(0.7),
          f'\u201c{text}\u201d', sz=16, c=C["dark"], al=PP_ALIGN.LEFT, ls=1.5)
        if author:
            T(s, ML + Inches(0.3), y + Inches(0.85), CW - Inches(0.3), Inches(0.3),
              f"\u2014 {author}", sz=SZ["body_sm"], c=C["gray"])
        return total_h
    elif style == "box":
        total_h = 1.2 if not author else 1.5
        R(s, ML, y, CW, Inches(total_h), f=C["light"], lc=color, lw=1.5)
        # 큰 따옴표
        T(s, ML + Inches(0.2), y + Inches(0.02), Inches(0.5), Inches(0.5),
          "\u201c", sz=36, c=color, b=True)
        T(s, ML + Inches(0.5), y + Inches(0.15), CW - Inches(0.8), Inches(0.7),
          text, sz=14, c=C["dark"], al=PP_ALIGN.LEFT, ls=1.5)
        if author:
            T(s, ML + Inches(0.5), y + Inches(0.9), CW - Inches(0.8), Inches(0.3),
              f"\u2014 {author}", sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.RIGHT)
        return total_h


def NUMBERED_LIST(s, x, y, w, items, sz=13, gap=0.55):
    """번호 리스트 — 색상 원형 번호 + 제목 + 설명

    Args:
        items: [("제목", "설명"), ...] 또는 ["항목1", "항목2", ...]
        gap: 항목 간 간격
    Returns:
        float: 전체 높이 (inches)
    """
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
    w_in = float(w / 914400) if hasattr(w, '__class__') and w.__class__.__name__ != 'float' else w
    total_h = 0
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            title, desc = item
        else:
            title, desc = item, ""
        iy = y + Inches(gap * i)
        clr = palette[i % len(palette)]
        # 번호 원 (라운드 박스)
        RBOX(s, x, iy, Inches(0.38), Inches(0.38), clr,
             str(i + 1), sz=12, tc=C["white"], b=True, radius=0.5)
        # 제목
        T(s, x + Inches(0.52), iy + Inches(0.02), Inches(w_in - 0.52), Inches(0.25),
          title, sz=sz, c=C["dark"], b=True)
        # 설명
        if desc:
            T(s, x + Inches(0.52), iy + Inches(0.28), Inches(w_in - 0.52), Inches(0.22),
              desc, sz=SZ["body_sm"], c=C["gray"])
        total_h = gap * i + gap
    return total_h


# ═══════════════════════════════════════════════════════════════
#  19. v3.5 — 네이티브 차트 (BAR / PIE / LINE)
# ═══════════════════════════════════════════════════════════════

def BAR_CHART(s, x, y, w, h, categories, series_data, title="",
              chart_type="column", colors=None):
    """바 차트 — 비교 데이터 시각화

    Args:
        categories: ["항목A", "항목B", "항목C"]
        series_data: [("시리즈명", [10, 20, 30]), ...]
        chart_type: "column" (세로) / "bar" (가로) / "stacked"
        colors: 시리즈별 색상 리스트
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked": XL_CHART_TYPE.COLUMN_STACKED,
    }
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    graphic = s.shapes.add_chart(type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                                  x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(SZ["caption"])
        chart.legend.font.name = FONT

    # Modern 스타일 적용
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    plot = chart.plots[0]
    plot.gap_width = 80
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]

    # 축 폰트
    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.category_axis.tick_labels.font.name = FONT
    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.has_major_gridlines = True

    return graphic


def PIE_CHART(s, x, y, w, h, categories, values, title="", colors=None, donut=False):
    """파이/도넛 차트

    Args:
        categories: ["항목A", "항목B", "항목C"]
        values: [30, 50, 20]
        donut: True면 도넛 차트
        colors: 항목별 색상 리스트
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("값", values)

    graphic = s.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(SZ["caption"])
    chart.legend.font.name = FONT

    # Modern 스타일 색상
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"],
                  C["green"], C["gold"], C["gray"], C["orange"]]
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]

    return graphic


def LINE_CHART(s, x, y, w, h, categories, series_data, title="",
               colors=None, smooth=False):
    """라인 차트 — 추세 데이터

    Args:
        categories: ["1월", "2월", "3월", ...]
        series_data: [("시리즈명", [10, 20, 30, ...]), ...]
        smooth: True면 곡선
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type = XL_CHART_TYPE.LINE_MARKERS if not smooth else XL_CHART_TYPE.LINE
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    graphic = s.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(SZ["caption"])
        chart.legend.font.name = FONT

    # Modern 스타일 적용
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    plot = chart.plots[0]
    if smooth:
        plot.smooth = True
    for i, series in enumerate(plot.series):
        series.format.line.color.rgb = colors[i % len(colors)]
        series.format.line.width = Pt(2.5)
        series.smooth = smooth

    # 축 폰트
    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.category_axis.tick_labels.font.name = FONT
    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.has_major_gridlines = True

    return graphic


# ═══════════════════════════════════════════════════════════════
#  20. v3.5 — 테마 시스템
# ═══════════════════════════════════════════════════════════════

THEMES = {
    "default_blue": {
        "primary": (0, 44, 95), "secondary": (0, 170, 210),
        "teal": (0, 161, 156), "accent": (230, 51, 18),
        "dark": (33, 33, 33), "light": (245, 245, 245),
    },
    "warm": {
        "primary": (139, 69, 19), "secondary": (210, 105, 30),
        "teal": (160, 82, 45), "accent": (220, 20, 60),
        "dark": (50, 30, 15), "light": (255, 248, 240),
    },
    "forest": {
        "primary": (27, 94, 32), "secondary": (76, 175, 80),
        "teal": (0, 150, 136), "accent": (255, 111, 0),
        "dark": (30, 40, 30), "light": (241, 248, 233),
    },
    "corporate": {
        "primary": (38, 50, 56), "secondary": (255, 111, 0),
        "teal": (0, 137, 123), "accent": (213, 0, 0),
        "dark": (33, 33, 33), "light": (245, 245, 245),
    },
    "purple": {
        "primary": (74, 20, 140), "secondary": (156, 39, 176),
        "teal": (0, 150, 136), "accent": (233, 30, 99),
        "dark": (40, 20, 60), "light": (243, 229, 245),
    },
}

_original_colors = {k: v for k, v in C.items()}


def apply_theme(theme_name):
    """컬러 팔레트 일괄 변경

    Args:
        theme_name: THEMES 딕셔너리의 키 (예: "default_blue", "warm", "forest")
    """
    theme = THEMES.get(theme_name)
    if not theme:
        raise ValueError(f"Unknown theme: {theme_name}. Available: {list(THEMES.keys())}")
    for key, rgb in theme.items():
        if key in C:
            C[key] = RGBColor(*rgb)
    return theme_name


def reset_theme():
    """테마를 기본값(default_blue)으로 복원"""
    for k, v in _original_colors.items():
        C[k] = v


def list_themes():
    """사용 가능한 테마 목록 출력"""
    for name in THEMES:
        print(f"  {name}")


# ═══════════════════════════════════════════════════════════════
#  21. v3.5 — 레이아웃 시퀀스 검증
# ═══════════════════════════════════════════════════════════════

def validate_sequence(slide_info):
    """레이아웃 시퀀스 검증 → 경고 메시지 리스트 반환

    Args:
        slide_info: [{"layout": "THREE_COL", "has_image": False, "has_highlight": True}, ...]

    Returns:
        list[str] — 경고 메시지
    """
    warnings = []
    layouts = [s.get("layout", "") for s in slide_info]

    # 규칙 1: 같은 레이아웃 3회 연속 금지
    for i in range(len(layouts) - 2):
        if layouts[i] and layouts[i] == layouts[i + 1] == layouts[i + 2]:
            warnings.append(
                f"[시각 단조] 슬라이드 {i+1}-{i+3}: '{layouts[i]}' 3회 연속 → 다른 레이아웃 권장")

    # 규칙 2: 5장 연속 이미지 없음 경고
    no_img_streak = 0
    for i, s in enumerate(slide_info):
        if s.get("has_image", False):
            no_img_streak = 0
        else:
            no_img_streak += 1
            if no_img_streak >= 5:
                warnings.append(
                    f"[이미지 부족] 슬라이드 {i-3}~{i+1}: 5장 연속 이미지 없음 → IMG_PH 추가 권장")
                no_img_streak = 0

    # 규칙 3: 5장 연속 HIGHLIGHT 없음 경고
    no_hl_streak = 0
    for i, s in enumerate(slide_info):
        if s.get("has_highlight", False):
            no_hl_streak = 0
        else:
            no_hl_streak += 1
            if no_hl_streak >= 5:
                warnings.append(
                    f"[강조 부족] 슬라이드 {i-3}~{i+1}: 5장 연속 HIGHLIGHT 없음 → 핵심 메시지 강조 권장")
                no_hl_streak = 0

    return warnings


# ═══════════════════════════════════════════════════════════════
#  22. v3.5 — new_presentation 템플릿 지원
# ═══════════════════════════════════════════════════════════════

def new_presentation_from_template(template_path):
    """기존 PPTX 템플릿 기반 프레젠테이션 생성

    마스터 슬라이드의 로고, 푸터, 테마 색상을 그대로 활용.
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    prs = Presentation(template_path)
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


# ═══════════════════════════════════════════════════════════════
#  23. v3.5 — 카드 컴포넌트 (그림자 + 라운드)
# ═══════════════════════════════════════════════════════════════

def CARD(s, x, y, w, h, title, body="", color=None, shadow=True, rounded=True):
    """고급 카드 — 라운드 + 미세 그림자 + 컬러 상단바

    Args:
        title: 카드 제목
        body: 본문 텍스트 또는 ["줄1", "줄2"]
        color: 상단 바 색상
        shadow: 그림자 여부
        rounded: 라운드 코너 여부
    """
    if color is None:
        color = C["primary"]
    h_in = float(h / 914400) if hasattr(h, '__class__') and h.__class__.__name__ != 'float' else h
    h = Inches(h_in) if not hasattr(h, '__class__') or h.__class__.__name__ == 'float' else h

    # 배경 카드
    if rounded:
        card_sh = RBOX(s, x, y, w, h, C["white"], radius=0.08)
        card_sh.line.color.rgb = C["lgray"]
        card_sh.line.width = Pt(0.5)
    else:
        card_sh = R(s, x, y, w, h, f=C["white"], lc=C["lgray"], lw=0.5)

    if shadow:
        add_shadow(card_sh)

    # 상단 컬러 바
    w_in = float(w / 914400) if hasattr(w, '__class__') and w.__class__.__name__ != 'float' else w
    R(s, x, y, w, Pt(5), f=color)

    # 제목
    T(s, x + Inches(0.15), y + Inches(0.15), Inches(w_in - 0.3), Inches(0.35),
      title, sz=SZ["body"], c=C["dark"], b=True)

    # 본문
    if body:
        body_y = y + Inches(0.55)
        body_h = Inches(h_in - 0.7)
        if isinstance(body, list):
            MT(s, x + Inches(0.15), body_y, Inches(w_in - 0.3), body_h,
               body, sz=SZ["body_sm"], bul=True)
        else:
            T(s, x + Inches(0.15), body_y, Inches(w_in - 0.3), body_h,
              body, sz=SZ["body_sm"], c=C["gray"], ls=1.5)

    return card_sh


# 별칭
_p = new_slide
_set_char_spacing = set_char_spacing
