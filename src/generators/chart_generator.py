"""
차트 생성기 ([회사명])

타임라인, 조직도, 차트 등 시각화 요소 생성
"""

from typing import Any, Dict, List, Optional

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .template_manager import TemplateManager
from ..schemas.proposal_schema import TimelineItem, OrgChartNode, KPIItem, CompetitorComparison
from ..utils.logger import get_logger

logger = get_logger("chart_generator")


class ChartGenerator:
    """차트/다이어그램 생성기"""

    def __init__(self, template_manager: TemplateManager):
        self.template_manager = template_manager
        self.design = template_manager.design_system

    def add_timeline_to_slide(
        self,
        slide,
        timeline_items: List[TimelineItem],
        top: float = 2.0,
    ) -> None:
        """
        슬라이드에 타임라인 추가

        Args:
            slide: 대상 슬라이드
            timeline_items: 타임라인 항목 목록
            top: 상단 여백 (인치)
        """
        if not timeline_items:
            return

        num_phases = len(timeline_items)
        total_width = 12.0
        phase_width = total_width / num_phases
        left_start = 0.67

        # 각 Phase 박스 추가
        colors = [
            self.template_manager.get_color("primary"),
            self.template_manager.get_color("secondary"),
            RGBColor(100, 149, 237),  # Cornflower blue
            RGBColor(60, 179, 113),  # Medium sea green
        ]

        for i, item in enumerate(timeline_items):
            left = left_start + (i * phase_width)
            color = colors[i % len(colors)]

            # Phase 박스
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(phase_width - 0.2),
                Inches(1.2),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            # Phase 텍스트
            tf = shape.text_frame
            tf.word_wrap = True

            # Phase 번호/이름
            p1 = tf.paragraphs[0]
            p1.text = item.phase
            p1.font.size = Pt(14)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            p1.alignment = PP_ALIGN.CENTER

            # Phase 제목
            p2 = tf.add_paragraph()
            p2.text = item.title
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER

            # 기간 텍스트 (박스 아래)
            duration_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top + 1.3),
                Inches(phase_width - 0.2),
                Inches(0.4),
            )
            duration_tf = duration_box.text_frame
            duration_p = duration_tf.paragraphs[0]
            duration_p.text = item.duration
            duration_p.font.size = Pt(11)
            duration_p.font.color.rgb = self.template_manager.get_color("text_light")
            duration_p.alignment = PP_ALIGN.CENTER

            # 설명 또는 마일스톤 (있는 경우)
            if item.description or item.milestones:
                desc_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top + 1.8),
                    Inches(phase_width - 0.2),
                    Inches(2.5),
                )
                desc_tf = desc_box.text_frame
                desc_tf.word_wrap = True

                if item.description:
                    desc_p = desc_tf.paragraphs[0]
                    desc_p.text = item.description
                    desc_p.font.size = Pt(10)

                if item.milestones:
                    for milestone in item.milestones:
                        ms_p = desc_tf.add_paragraph()
                        ms_p.text = f"• {milestone}"
                        ms_p.font.size = Pt(10)

        # 화살표 추가 (Phase 간)
        arrow_y = top + 0.5
        for i in range(num_phases - 1):
            left = left_start + ((i + 1) * phase_width) - 0.25

            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(left),
                Inches(arrow_y),
                Inches(0.3),
                Inches(0.2),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = self.template_manager.get_color("text_light")
            arrow.line.fill.background()

    def add_org_chart_to_slide(
        self,
        slide,
        org_chart: OrgChartNode,
        top: float = 1.8,
    ) -> None:
        """
        슬라이드에 조직도 추가

        Args:
            slide: 대상 슬라이드
            org_chart: 조직도 루트 노드
            top: 상단 여백 (인치)
        """
        if not org_chart:
            return

        # 루트 노드 (PM)
        root_left = 5.5
        root_width = 2.5
        root_height = 0.9

        self._add_org_node(
            slide,
            org_chart.name,
            org_chart.role,
            root_left,
            top,
            root_width,
            root_height,
            is_root=True,
        )

        # 자식 노드
        if org_chart.children:
            num_children = len(org_chart.children)
            child_width = 2.2
            child_height = 0.8
            child_top = top + 1.5

            # 총 너비 계산
            total_width = num_children * child_width + (num_children - 1) * 0.3
            start_left = (13.33 - total_width) / 2

            for i, child in enumerate(org_chart.children):
                child_left = start_left + i * (child_width + 0.3)

                # 연결선
                self._add_connector(
                    slide,
                    root_left + root_width / 2,
                    top + root_height,
                    child_left + child_width / 2,
                    child_top,
                )

                # 자식 노드
                self._add_org_node(
                    slide,
                    child.name,
                    child.role,
                    child_left,
                    child_top,
                    child_width,
                    child_height,
                    is_root=False,
                )

                # 손자 노드 (2단계까지만)
                if child.children:
                    grandchild_top = child_top + 1.3
                    gc_width = 2.0
                    gc_height = 0.7

                    num_gc = len(child.children)
                    gc_total_width = num_gc * gc_width + (num_gc - 1) * 0.2
                    gc_start_left = child_left + (child_width - gc_total_width) / 2

                    for j, grandchild in enumerate(child.children[:3]):  # 최대 3개
                        gc_left = gc_start_left + j * (gc_width + 0.2)

                        self._add_connector(
                            slide,
                            child_left + child_width / 2,
                            child_top + child_height,
                            gc_left + gc_width / 2,
                            grandchild_top,
                        )

                        self._add_org_node(
                            slide,
                            grandchild.name,
                            grandchild.role,
                            gc_left,
                            grandchild_top,
                            gc_width,
                            gc_height,
                            is_root=False,
                        )

    def _add_org_node(
        self,
        slide,
        name: str,
        role: str,
        left: float,
        top: float,
        width: float,
        height: float,
        is_root: bool = False,
    ) -> None:
        """조직도 노드 추가"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )

        if is_root:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.template_manager.get_color("primary")
            text_color = RGBColor(255, 255, 255)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.template_manager.get_color(
                "background_light"
            )
            text_color = self.template_manager.get_color("text_dark")

        shape.line.color.rgb = self.template_manager.get_color("primary")
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True

        # 이름
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(11) if is_root else Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = text_color
        p1.alignment = PP_ALIGN.CENTER

        # 역할
        p2 = tf.add_paragraph()
        p2.text = role
        p2.font.size = Pt(9)
        p2.font.color.rgb = text_color if is_root else self.template_manager.get_color(
            "text_light"
        )
        p2.alignment = PP_ALIGN.CENTER

    def _add_connector(
        self,
        slide,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
    ) -> None:
        """연결선 추가"""
        # 수직선 (시작점에서 아래로)
        mid_y = (start_y + end_y) / 2

        line1 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x),
            Inches(start_y),
            Inches(start_x),
            Inches(mid_y),
        )
        line1.line.color.rgb = self.template_manager.get_color("text_light")
        line1.line.width = Pt(1)

        # 수평선
        line2 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x),
            Inches(mid_y),
            Inches(end_x),
            Inches(mid_y),
        )
        line2.line.color.rgb = self.template_manager.get_color("text_light")
        line2.line.width = Pt(1)

        # 수직선 (끝점으로)
        line3 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(end_x),
            Inches(mid_y),
            Inches(end_x),
            Inches(end_y),
        )
        line3.line.color.rgb = self.template_manager.get_color("text_light")
        line3.line.width = Pt(1)

    def add_simple_bar_chart_placeholder(
        self,
        slide,
        title: str,
        data: Dict[str, Any],
        left: float = 1.0,
        top: float = 2.0,
        width: float = 11.0,
        height: float = 4.5,
    ) -> None:
        """
        간단한 바 차트 플레이스홀더 (이미지 대용)

        실제 차트는 matplotlib으로 생성 후 이미지로 삽입 가능
        여기서는 텍스트 기반 표현 제공
        """
        # 차트 영역 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        shape.line.color.rgb = self.template_manager.get_color("text_light")

        # 차트 제목
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.2),
            Inches(top + 0.2),
            Inches(width - 0.4),
            Inches(0.5),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(14)
        title_p.font.bold = True

        # 데이터 텍스트 표현
        if "items" in data:
            items = data["items"]
            item_height = (height - 1.0) / max(len(items), 1)

            for i, item in enumerate(items):
                item_box = slide.shapes.add_textbox(
                    Inches(left + 0.5),
                    Inches(top + 0.8 + i * item_height),
                    Inches(width - 1.0),
                    Inches(item_height),
                )
                item_tf = item_box.text_frame
                item_p = item_tf.paragraphs[0]
                item_p.text = f"• {item.get('label', '')}: {item.get('value', '')}"
                item_p.font.size = Pt(12)

    # ==========================================================================
    # KPI 시각화
    # ==========================================================================
    def add_kpi_cards_to_slide(
        self,
        slide,
        kpis: List[KPIItem],
        top: float = 2.0,
    ) -> None:
        """
        KPI 카드 시각화 추가

        Args:
            slide: 대상 슬라이드
            kpis: KPIItem 목록
            top: 상단 여백
        """
        if not kpis:
            return

        num_kpis = min(len(kpis), 4)
        total_width = 12.0
        card_width = total_width / num_kpis
        left_start = 0.67

        for i, kpi in enumerate(kpis[:4]):
            left = left_start + (i * card_width)
            self._add_kpi_card(slide, kpi, left, top, card_width - 0.3)

    def _add_kpi_card(
        self,
        slide,
        kpi: KPIItem,
        left: float,
        top: float,
        width: float,
    ) -> None:
        """개별 KPI 카드"""
        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(2.8),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = self.template_manager.get_color("primary")
        card.line.width = Pt(2)

        # 지표명
        metric_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.2),
            Inches(width - 0.2),
            Inches(0.5),
        )
        metric_tf = metric_box.text_frame
        metric_p = metric_tf.paragraphs[0]
        metric_p.text = kpi.metric
        metric_p.font.size = Pt(12)
        metric_p.font.color.rgb = self.template_manager.get_color("text_light")
        metric_p.alignment = PP_ALIGN.CENTER

        # 목표값 (큰 숫자)
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.7),
            Inches(width - 0.2),
            Inches(1.0),
        )
        value_tf = value_box.text_frame
        value_p = value_tf.paragraphs[0]
        value_p.text = kpi.target
        value_p.font.size = Pt(36)
        value_p.font.bold = True
        value_p.font.color.rgb = self.template_manager.get_color("primary")
        value_p.alignment = PP_ALIGN.CENTER

        # 개선율
        if kpi.improvement:
            imp_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 1.7),
                Inches(width - 0.2),
                Inches(0.5),
            )
            imp_tf = imp_box.text_frame
            imp_p = imp_tf.paragraphs[0]
            imp_p.text = f"▲ {kpi.improvement}"
            imp_p.font.size = Pt(14)
            imp_p.font.bold = True
            imp_p.font.color.rgb = RGBColor(40, 167, 69)  # Green
            imp_p.alignment = PP_ALIGN.CENTER

        # 기준값
        if kpi.baseline:
            baseline_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 2.2),
                Inches(width - 0.2),
                Inches(0.4),
            )
            baseline_tf = baseline_box.text_frame
            baseline_p = baseline_tf.paragraphs[0]
            baseline_p.text = f"현재: {kpi.baseline}"
            baseline_p.font.size = Pt(10)
            baseline_p.font.color.rgb = self.template_manager.get_color("text_light")
            baseline_p.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # 경쟁사 비교 시각화
    # ==========================================================================
    def add_competitor_comparison_to_slide(
        self,
        slide,
        comparisons: List[CompetitorComparison],
        top: float = 1.8,
    ) -> None:
        """
        경쟁사 비교 시각화 추가

        Args:
            slide: 대상 슬라이드
            comparisons: CompetitorComparison 목록
            top: 상단 여백
        """
        if not comparisons:
            return

        # 헤더
        headers_y = top
        self._add_comparison_header(slide, "평가 기준", 0.5, headers_y, 2.5)
        self._add_comparison_header(
            slide, "당사 강점", 3.2, headers_y, 4.5,
            self.template_manager.get_color("primary")
        )
        self._add_comparison_header(
            slide, "경쟁사", 7.9, headers_y, 4.5,
            RGBColor(150, 150, 150)
        )

        # 각 비교 항목
        row_height = 0.9
        for i, comp in enumerate(comparisons[:5]):  # 최대 5개
            row_top = top + 0.8 + i * row_height
            self._add_comparison_row(slide, comp, row_top, row_height)

    def _add_comparison_header(
        self,
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        color: Optional[RGBColor] = None,
    ) -> None:
        """비교표 헤더"""
        if color is None:
            color = self.template_manager.get_color("text_dark")

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.6),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_comparison_row(
        self,
        slide,
        comp: CompetitorComparison,
        top: float,
        height: float,
    ) -> None:
        """비교표 행"""
        # 평가 기준
        criteria_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(top),
            Inches(2.5),
            Inches(height),
        )
        criteria_tf = criteria_box.text_frame
        criteria_tf.word_wrap = True
        criteria_p = criteria_tf.paragraphs[0]
        criteria_p.text = comp.criteria
        criteria_p.font.size = Pt(11)
        criteria_p.font.bold = True
        criteria_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 당사 강점 (강조)
        our_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.2),
            Inches(top),
            Inches(4.5),
            Inches(height - 0.1),
        )
        our_shape.fill.solid()
        our_shape.fill.fore_color.rgb = RGBColor(232, 245, 233)  # Light green
        our_shape.line.color.rgb = self.template_manager.get_color("primary")
        our_shape.line.width = Pt(1)

        our_tf = our_shape.text_frame
        our_tf.word_wrap = True
        our_tf.margin_left = Inches(0.1)
        our_p = our_tf.paragraphs[0]
        our_p.text = f"✓ {comp.our_strength}"
        our_p.font.size = Pt(11)
        our_p.font.color.rgb = self.template_manager.get_color("text_dark")
        our_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 경쟁사
        comp_box = slide.shapes.add_textbox(
            Inches(7.9),
            Inches(top),
            Inches(4.5),
            Inches(height),
        )
        comp_tf = comp_box.text_frame
        comp_tf.word_wrap = True
        comp_p = comp_tf.paragraphs[0]
        comp_p.text = comp.competitor
        comp_p.font.size = Pt(10)
        comp_p.font.color.rgb = RGBColor(120, 120, 120)
        comp_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==========================================================================
    # ROI 시각화
    # ==========================================================================
    def add_roi_visualization(
        self,
        slide,
        investment: str,
        annual_benefit: str,
        payback_period: str,
        roi_percent: str,
        top: float = 2.0,
    ) -> None:
        """
        ROI 시각화 추가

        Args:
            slide: 대상 슬라이드
            investment: 투자 비용
            annual_benefit: 연간 절감/이익
            payback_period: 투자 회수 기간
            roi_percent: ROI 퍼센트
        """
        center_x = 6.67

        # ROI 큰 원
        roi_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - 1.5),
            Inches(top),
            Inches(3.0),
            Inches(3.0),
        )
        roi_circle.fill.solid()
        roi_circle.fill.fore_color.rgb = self.template_manager.get_color("primary")
        roi_circle.line.fill.background()

        roi_tf = roi_circle.text_frame
        roi_tf.word_wrap = True
        roi_p1 = roi_tf.paragraphs[0]
        roi_p1.text = "ROI"
        roi_p1.font.size = Pt(16)
        roi_p1.font.color.rgb = RGBColor(255, 255, 255)
        roi_p1.alignment = PP_ALIGN.CENTER

        roi_p2 = roi_tf.add_paragraph()
        roi_p2.text = roi_percent
        roi_p2.font.size = Pt(40)
        roi_p2.font.bold = True
        roi_p2.font.color.rgb = RGBColor(255, 255, 255)
        roi_p2.alignment = PP_ALIGN.CENTER

        roi_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 세부 항목들 (아래)
        items = [
            ("투자 비용", investment),
            ("연간 절감 효과", annual_benefit),
            ("투자 회수 기간", payback_period),
        ]

        item_width = 3.5
        items_left = (13.33 - len(items) * item_width) / 2
        items_top = top + 3.5

        for i, (label, value) in enumerate(items):
            item_left = items_left + i * item_width

            # 레이블
            label_box = slide.shapes.add_textbox(
                Inches(item_left),
                Inches(items_top),
                Inches(item_width - 0.3),
                Inches(0.4),
            )
            label_tf = label_box.text_frame
            label_p = label_tf.paragraphs[0]
            label_p.text = label
            label_p.font.size = Pt(11)
            label_p.font.color.rgb = self.template_manager.get_color("text_light")
            label_p.alignment = PP_ALIGN.CENTER

            # 값
            value_box = slide.shapes.add_textbox(
                Inches(item_left),
                Inches(items_top + 0.4),
                Inches(item_width - 0.3),
                Inches(0.6),
            )
            value_tf = value_box.text_frame
            value_p = value_tf.paragraphs[0]
            value_p.text = value
            value_p.font.size = Pt(18)
            value_p.font.bold = True
            value_p.font.color.rgb = self.template_manager.get_color("primary")
            value_p.alignment = PP_ALIGN.CENTER
