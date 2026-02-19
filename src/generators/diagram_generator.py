"""
다이어그램 생성기 ([회사명])

제안서 내용에 맞는 도식화 요소 생성:
- 프로세스 플로우
- 피처 박스 (3-4열)
- 아이콘 카드
- KPI 대시보드
- Before/After 비교
- 컨셉 다이어그램
"""

from typing import Any, Dict, List, Optional, Tuple

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .template_manager import TemplateManager
from ..utils.logger import get_logger

logger = get_logger("diagram_generator")


class DiagramGenerator:
    """도식화 요소 생성기"""

    def __init__(self, template_manager: TemplateManager):
        self.template_manager = template_manager
        self.design = template_manager.design_system

    # ==========================================================================
    # 프로세스 플로우 다이어그램
    # ==========================================================================
    def add_process_flow(
        self,
        slide,
        steps: List[Dict[str, str]],
        top: float = 2.5,
        style: str = "arrow",  # arrow, chevron, circle
    ) -> None:
        """
        프로세스 플로우 다이어그램 추가

        Args:
            slide: 대상 슬라이드
            steps: [{"title": "단계1", "description": "설명"}, ...]
            top: 상단 여백
            style: 스타일 (arrow, chevron, circle)
        """
        if not steps:
            return

        num_steps = len(steps)
        total_width = 12.0
        step_width = total_width / num_steps
        left_start = 0.67

        colors = self._get_gradient_colors(num_steps)

        for i, step in enumerate(steps):
            left = left_start + (i * step_width)
            color = colors[i]

            if style == "chevron":
                self._add_chevron_step(slide, step, left, top, step_width - 0.15, color, i == 0)
            elif style == "circle":
                self._add_circle_step(slide, step, left, top, step_width - 0.2, color, i + 1)
            else:  # arrow (default)
                self._add_arrow_step(slide, step, left, top, step_width - 0.2, color)

            # 화살표 연결 (마지막 제외)
            if i < num_steps - 1 and style == "arrow":
                self._add_flow_arrow(slide, left + step_width - 0.25, top + 0.5)

    def _add_arrow_step(
        self, slide, step: Dict, left: float, top: float, width: float, color: RGBColor
    ) -> None:
        """화살표 스타일 단계 박스"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)

        # 제목
        p1 = tf.paragraphs[0]
        p1.text = step.get("title", "")
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER

        # 설명
        if step.get("description"):
            p2 = tf.add_paragraph()
            p2.text = step["description"]
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER

    def _add_chevron_step(
        self, slide, step: Dict, left: float, top: float, width: float, color: RGBColor, is_first: bool
    ) -> None:
        """쉐브론(펜타곤) 스타일 단계"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON if not is_first else MSO_SHAPE.PENTAGON,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = step.get("title", "")
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER

    def _add_circle_step(
        self, slide, step: Dict, left: float, top: float, width: float, color: RGBColor, number: int
    ) -> None:
        """원형 스타일 단계 (번호 포함)"""
        center_left = left + (width - 1.2) / 2

        # 원형 배경
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_left),
            Inches(top),
            Inches(1.2),
            Inches(1.2),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # 번호
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 제목 (아래)
        title_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top + 1.4),
            Inches(width),
            Inches(0.8),
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = step.get("title", "")
        title_p.font.size = Pt(11)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER

    def _add_flow_arrow(self, slide, left: float, top: float) -> None:
        """플로우 화살표"""
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(left),
            Inches(top),
            Inches(0.3),
            Inches(0.25),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.template_manager.get_color("text_light")
        arrow.line.fill.background()

    # ==========================================================================
    # 피처 박스 (3-4열 그리드)
    # ==========================================================================
    def add_feature_boxes(
        self,
        slide,
        features: List[Dict[str, str]],
        top: float = 2.0,
        columns: int = 3,
        with_icons: bool = True,
    ) -> None:
        """
        피처 박스 그리드 추가

        Args:
            slide: 대상 슬라이드
            features: [{"title": "기능1", "description": "설명", "icon": "★"}, ...]
            top: 상단 여백
            columns: 열 개수 (3 또는 4)
            with_icons: 아이콘 표시 여부
        """
        if not features:
            return

        total_width = 12.0
        margin = 0.3
        box_width = (total_width - margin * (columns + 1)) / columns
        box_height = 1.8
        left_start = 0.67

        rows = (len(features) + columns - 1) // columns

        for i, feature in enumerate(features):
            row = i // columns
            col = i % columns

            left = left_start + margin + col * (box_width + margin)
            current_top = top + row * (box_height + margin)

            self._add_feature_box(
                slide, feature, left, current_top, box_width, box_height, with_icons
            )

    def _add_feature_box(
        self,
        slide,
        feature: Dict,
        left: float,
        top: float,
        width: float,
        height: float,
        with_icons: bool,
    ) -> None:
        """개별 피처 박스"""
        # 박스 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        shape.line.color.rgb = self.template_manager.get_color("primary")
        shape.line.width = Pt(1.5)

        content_left = left + 0.15
        content_width = width - 0.3
        current_top = top + 0.15

        # 아이콘 (있는 경우)
        if with_icons and feature.get("icon"):
            icon_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top),
                Inches(content_width),
                Inches(0.5),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = feature["icon"]
            icon_p.font.size = Pt(24)
            icon_p.alignment = PP_ALIGN.CENTER
            current_top += 0.5

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(content_left),
            Inches(current_top),
            Inches(content_width),
            Inches(0.5),
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = feature.get("title", "")
        title_p.font.size = Pt(12)
        title_p.font.bold = True
        title_p.font.color.rgb = self.template_manager.get_color("primary")
        title_p.alignment = PP_ALIGN.CENTER

        # 설명
        if feature.get("description"):
            desc_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top + 0.45),
                Inches(content_width),
                Inches(height - current_top - 0.3),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = feature["description"]
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = self.template_manager.get_color("text_dark")
            desc_p.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # KPI 대시보드
    # ==========================================================================
    def add_kpi_dashboard(
        self,
        slide,
        kpis: List[Dict[str, Any]],
        top: float = 2.0,
    ) -> None:
        """
        KPI 대시보드 추가 (숫자 강조)

        Args:
            slide: 대상 슬라이드
            kpis: [{"metric": "만족도", "value": "95%", "baseline": "70%", "improvement": "+25%p"}, ...]
            top: 상단 여백
        """
        if not kpis:
            return

        num_kpis = min(len(kpis), 4)  # 최대 4개
        total_width = 12.0
        kpi_width = total_width / num_kpis
        left_start = 0.67

        for i, kpi in enumerate(kpis[:4]):
            left = left_start + (i * kpi_width)
            self._add_kpi_card(slide, kpi, left, top, kpi_width - 0.3)

    def _add_kpi_card(
        self, slide, kpi: Dict, left: float, top: float, width: float
    ) -> None:
        """개별 KPI 카드"""
        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(2.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = self.template_manager.get_color("primary")
        card.line.width = Pt(2)

        # 지표명
        metric_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.2),
            Inches(width - 0.2),
            Inches(0.5),
        )
        metric_tf = metric_box.text_frame
        metric_p = metric_tf.paragraphs[0]
        metric_p.text = kpi.get("metric", "")
        metric_p.font.size = Pt(12)
        metric_p.font.color.rgb = self.template_manager.get_color("text_light")
        metric_p.alignment = PP_ALIGN.CENTER

        # 목표값 (큰 숫자)
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.7),
            Inches(width - 0.2),
            Inches(1.0),
        )
        value_tf = value_box.text_frame
        value_p = value_tf.paragraphs[0]
        value_p.text = kpi.get("value", kpi.get("target", ""))
        value_p.font.size = Pt(36)
        value_p.font.bold = True
        value_p.font.color.rgb = self.template_manager.get_color("primary")
        value_p.alignment = PP_ALIGN.CENTER

        # 개선율 (있는 경우)
        improvement = kpi.get("improvement")
        if improvement:
            imp_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 1.7),
                Inches(width - 0.2),
                Inches(0.5),
            )
            imp_tf = imp_box.text_frame
            imp_p = imp_tf.paragraphs[0]
            imp_p.text = f"▲ {improvement}"
            imp_p.font.size = Pt(14)
            imp_p.font.bold = True
            # 개선이 양수면 녹색, 아니면 빨간색
            if improvement.startswith("-"):
                imp_p.font.color.rgb = RGBColor(220, 53, 69)
            else:
                imp_p.font.color.rgb = RGBColor(40, 167, 69)
            imp_p.alignment = PP_ALIGN.CENTER

        # 기준값 (있는 경우)
        baseline = kpi.get("baseline")
        if baseline:
            baseline_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 2.1),
                Inches(width - 0.2),
                Inches(0.3),
            )
            baseline_tf = baseline_box.text_frame
            baseline_p = baseline_tf.paragraphs[0]
            baseline_p.text = f"현재: {baseline}"
            baseline_p.font.size = Pt(10)
            baseline_p.font.color.rgb = self.template_manager.get_color("text_light")
            baseline_p.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # Before/After 비교 다이어그램
    # ==========================================================================
    def add_before_after(
        self,
        slide,
        items: List[Dict[str, str]],
        before_title: str = "AS-IS",
        after_title: str = "TO-BE",
        top: float = 1.8,
    ) -> None:
        """
        Before/After 비교 다이어그램

        Args:
            slide: 대상 슬라이드
            items: [{"label": "처리시간", "before": "30분", "after": "10분"}, ...]
            before_title: Before 제목
            after_title: After 제목
            top: 상단 여백
        """
        if not items:
            return

        left_col = 1.0
        right_col = 7.0
        col_width = 5.5
        arrow_left = 6.0

        # Before 헤더
        self._add_comparison_header(
            slide, before_title, left_col, top, col_width,
            RGBColor(108, 117, 125)  # 회색
        )

        # After 헤더
        self._add_comparison_header(
            slide, after_title, right_col, top, col_width,
            self.template_manager.get_color("primary")
        )

        # 화살표
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(arrow_left),
            Inches(top + 0.15),
            Inches(0.8),
            Inches(0.5),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.template_manager.get_color("accent")
        arrow.line.fill.background()

        # 항목들
        item_top = top + 1.0
        item_height = 0.8

        for i, item in enumerate(items[:5]):  # 최대 5개
            current_top = item_top + i * item_height

            # Before 값
            self._add_comparison_item(
                slide, item.get("label", ""), item.get("before", ""),
                left_col, current_top, col_width, is_before=True
            )

            # After 값
            self._add_comparison_item(
                slide, item.get("label", ""), item.get("after", ""),
                right_col, current_top, col_width, is_before=False
            )

    def _add_comparison_header(
        self, slide, title: str, left: float, top: float, width: float, color: RGBColor
    ) -> None:
        """비교 헤더"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.7),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_comparison_item(
        self, slide, label: str, value: str, left: float, top: float, width: float, is_before: bool
    ) -> None:
        """비교 항목"""
        # 레이블
        label_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width * 0.5),
            Inches(0.6),
        )
        label_tf = label_box.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.text = label
        label_p.font.size = Pt(11)
        label_p.font.color.rgb = self.template_manager.get_color("text_dark")

        # 값
        value_box = slide.shapes.add_textbox(
            Inches(left + width * 0.5),
            Inches(top),
            Inches(width * 0.5),
            Inches(0.6),
        )
        value_tf = value_box.text_frame
        value_p = value_tf.paragraphs[0]
        value_p.text = value
        value_p.font.size = Pt(14)
        value_p.font.bold = True
        value_p.alignment = PP_ALIGN.RIGHT

        if is_before:
            value_p.font.color.rgb = RGBColor(108, 117, 125)
        else:
            value_p.font.color.rgb = self.template_manager.get_color("primary")

    # ==========================================================================
    # 컨셉 다이어그램 (중앙 + 주변)
    # ==========================================================================
    def add_concept_diagram(
        self,
        slide,
        center: Dict[str, str],
        surrounding: List[Dict[str, str]],
        top: float = 2.0,
    ) -> None:
        """
        컨셉 다이어그램 (중앙에 핵심 컨셉, 주변에 관련 요소)

        Args:
            slide: 대상 슬라이드
            center: {"title": "핵심 컨셉", "description": "설명"}
            surrounding: [{"title": "요소1", "description": "설명"}, ...]
            top: 상단 여백
        """
        center_x = 6.67
        center_y = top + 1.5
        center_size = 2.5

        # 중앙 원
        center_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - center_size / 2),
            Inches(center_y - center_size / 2 + 0.5),
            Inches(center_size),
            Inches(center_size),
        )
        center_shape.fill.solid()
        center_shape.fill.fore_color.rgb = self.template_manager.get_color("primary")
        center_shape.line.fill.background()

        # 중앙 텍스트
        tf = center_shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = center.get("title", "")
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 주변 요소 배치 (원형 배열)
        num_items = min(len(surrounding), 6)
        if num_items == 0:
            return

        import math
        radius = 3.0
        angle_step = 360 / num_items

        for i, item in enumerate(surrounding[:6]):
            angle = math.radians(90 + i * angle_step)  # 12시 방향부터 시작
            item_x = center_x + radius * math.cos(angle)
            item_y = center_y + 0.5 + radius * math.sin(angle)

            # 연결선
            self._add_concept_connector(slide, center_x, center_y + 0.5, item_x, item_y)

            # 주변 박스
            self._add_concept_box(slide, item, item_x, item_y)

    def _add_concept_connector(
        self, slide, start_x: float, start_y: float, end_x: float, end_y: float
    ) -> None:
        """컨셉 연결선"""
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x),
            Inches(start_y),
            Inches(end_x),
            Inches(end_y),
        )
        connector.line.color.rgb = self.template_manager.get_color("text_light")
        connector.line.width = Pt(1)

    def _add_concept_box(self, slide, item: Dict, x: float, y: float) -> None:
        """컨셉 주변 박스"""
        box_width = 2.0
        box_height = 1.0

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - box_width / 2),
            Inches(y - box_height / 2),
            Inches(box_width),
            Inches(box_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        shape.line.color.rgb = self.template_manager.get_color("secondary")
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("title", "")
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.template_manager.get_color("text_dark")
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==========================================================================
    # 경쟁 비교표 (시각화)
    # ==========================================================================
    def add_competitive_comparison(
        self,
        slide,
        criteria: List[str],
        our_company: Dict[str, Any],
        competitors: List[Dict[str, Any]],
        top: float = 1.8,
    ) -> None:
        """
        경쟁 비교 시각화

        Args:
            slide: 대상 슬라이드
            criteria: ["기술력", "가격", "실적"]
            our_company: {"name": "당사", "scores": [5, 4, 5]}
            competitors: [{"name": "A사", "scores": [3, 5, 3]}, ...]
            top: 상단 여백
        """
        if not criteria or not our_company:
            return

        num_criteria = len(criteria)
        left_start = 1.5
        row_height = 0.8
        bar_width = 8.0

        # 헤더
        header_box = slide.shapes.add_textbox(
            Inches(left_start),
            Inches(top),
            Inches(2.0),
            Inches(0.5),
        )
        header_tf = header_box.text_frame
        header_p = header_tf.paragraphs[0]
        header_p.text = "평가 항목"
        header_p.font.size = Pt(11)
        header_p.font.bold = True

        # 각 기준별 바 차트
        for i, criterion in enumerate(criteria):
            current_top = top + 0.7 + i * row_height

            # 기준명
            label_box = slide.shapes.add_textbox(
                Inches(left_start),
                Inches(current_top),
                Inches(2.0),
                Inches(row_height),
            )
            label_tf = label_box.text_frame
            label_p = label_tf.paragraphs[0]
            label_p.text = criterion
            label_p.font.size = Pt(11)

            # 당사 바
            our_score = our_company.get("scores", [])[i] if i < len(our_company.get("scores", [])) else 3
            self._add_comparison_bar(
                slide, 3.5, current_top + 0.1, bar_width * (our_score / 5),
                self.template_manager.get_color("primary"), our_company.get("name", "당사")
            )

            # 경쟁사 바 (첫 번째만)
            if competitors:
                comp = competitors[0]
                comp_score = comp.get("scores", [])[i] if i < len(comp.get("scores", [])) else 3
                self._add_comparison_bar(
                    slide, 3.5, current_top + 0.4, bar_width * (comp_score / 5),
                    RGBColor(180, 180, 180), comp.get("name", "경쟁사")
                )

    def _add_comparison_bar(
        self, slide, left: float, top: float, width: float, color: RGBColor, label: str
    ) -> None:
        """비교 바"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(max(width, 0.3)),
            Inches(0.25),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    # ==========================================================================
    # 유틸리티
    # ==========================================================================
    def _get_gradient_colors(self, count: int) -> List[RGBColor]:
        """그라데이션 색상 목록 생성"""
        primary = self.template_manager.get_color("primary")
        secondary = self.template_manager.get_color("secondary")

        colors = [
            primary,
            secondary,
            RGBColor(100, 149, 237),  # Cornflower blue
            RGBColor(60, 179, 113),   # Medium sea green
            RGBColor(147, 112, 219),  # Medium purple
            RGBColor(255, 127, 80),   # Coral
        ]

        return colors[:count] if count <= len(colors) else colors * (count // len(colors) + 1)
