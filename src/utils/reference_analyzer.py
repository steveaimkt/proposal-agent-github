"""
레퍼런스 PPTX 분석기 (v1.0)

레퍼런스 제안서 PPTX 파일에서 디자인 요소를 추출하여
slide_kit 테마 및 디자인 시스템에 반영합니다.

추출 항목:
- 컬러 팔레트 (배경색, 텍스트색, 강조색)
- 폰트 정보 (폰트명, 크기, 굵기)
- 레이아웃 패턴 (슬라이드별 요소 배치)
- 슬라이드 구조 (섹션 구분, 콘텐츠 유형)
"""

import json
from collections import Counter
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.dml.color import RGBColor

from ..utils.logger import get_logger

logger = get_logger("reference_analyzer")


class ReferenceAnalyzer:
    """레퍼런스 PPTX에서 디자인 요소를 추출하는 분석기"""

    def __init__(self, pptx_path: Path):
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"레퍼런스 파일을 찾을 수 없습니다: {pptx_path}")
        self.prs = Presentation(str(self.pptx_path))
        self._colors: List[Tuple[int, int, int]] = []
        self._bg_colors: List[Tuple[int, int, int]] = []
        self._fonts: List[str] = []
        self._font_sizes: List[float] = []
        self._bold_usage: Counter = Counter()

    def analyze(self) -> Dict[str, Any]:
        """전체 분석 실행 → 디자인 프로파일 반환"""
        logger.info(f"레퍼런스 분석 시작: {self.pptx_path.name}")

        result = {
            "file_name": self.pptx_path.name,
            "slide_count": len(self.prs.slides),
            "slide_size": {
                "width": self.prs.slide_width,
                "height": self.prs.slide_height,
                "width_inches": round(self.prs.slide_width / 914400, 2),
                "height_inches": round(self.prs.slide_height / 914400, 2),
            },
            "colors": self._extract_colors(),
            "fonts": self._extract_fonts(),
            "layouts": self._extract_layouts(),
            "structure": self._extract_structure(),
        }

        logger.info(
            f"레퍼런스 분석 완료: {result['slide_count']}장, "
            f"주요 컬러 {len(result['colors'].get('primary_candidates', []))}개, "
            f"폰트 {len(result['fonts'].get('font_families', []))}종"
        )

        return result

    def _extract_colors(self) -> Dict[str, Any]:
        """슬라이드에서 사용된 컬러 추출"""
        text_colors = []
        bg_colors = []
        shape_fills = []

        for slide in self.prs.slides:
            # 배경색 추출
            bg = self._get_slide_bg_color(slide)
            if bg:
                bg_colors.append(bg)

            # 도형/텍스트 컬러 추출
            for shape in slide.shapes:
                # 도형 채우기 색상
                fill_color = self._get_shape_fill_color(shape)
                if fill_color:
                    shape_fills.append(fill_color)

                # 텍스트 색상
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    rgb = run.font.color.rgb
                                    if rgb:
                                        text_colors.append((rgb[0], rgb[1], rgb[2]))
                            except (AttributeError, TypeError):
                                pass

        # 빈도 기반 상위 컬러 추출
        all_colors = text_colors + shape_fills
        color_counter = Counter(all_colors)
        bg_counter = Counter(bg_colors)

        # 흰/검 제외한 의미있는 컬러 추출
        significant_colors = [
            (color, count) for color, count in color_counter.most_common(30)
            if not self._is_neutral(color)
        ]

        # 밝기 기준 분류
        dark_colors = [(c, n) for c, n in significant_colors if self._brightness(c) < 100]
        bright_colors = [(c, n) for c, n in significant_colors if self._brightness(c) >= 100]

        return {
            "primary_candidates": [
                {"rgb": c, "hex": self._rgb_to_hex(c), "count": n}
                for c, n in dark_colors[:5]
            ],
            "secondary_candidates": [
                {"rgb": c, "hex": self._rgb_to_hex(c), "count": n}
                for c, n in bright_colors[:5]
            ],
            "background_colors": [
                {"rgb": c, "hex": self._rgb_to_hex(c), "count": n}
                for c, n in bg_counter.most_common(5)
            ],
            "text_colors": [
                {"rgb": c, "hex": self._rgb_to_hex(c), "count": n}
                for c, n in Counter(text_colors).most_common(5)
            ],
            "total_unique_colors": len(set(all_colors)),
        }

    def _extract_fonts(self) -> Dict[str, Any]:
        """폰트 정보 추출"""
        font_names = []
        font_sizes = []
        bold_count = 0
        total_runs = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        total_runs += 1
                        if run.font.name:
                            font_names.append(run.font.name)
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
                        if run.font.bold:
                            bold_count += 1

        name_counter = Counter(font_names)
        size_counter = Counter(font_sizes)

        # 크기별 분류 (타이틀 vs 본문)
        sorted_sizes = sorted(size_counter.items(), key=lambda x: -x[1])
        title_sizes = [s for s, _ in sorted_sizes if s >= 24]
        body_sizes = [s for s, _ in sorted_sizes if 12 <= s < 24]

        return {
            "font_families": [
                {"name": name, "count": count}
                for name, count in name_counter.most_common(5)
            ],
            "primary_font": name_counter.most_common(1)[0][0] if name_counter else "Pretendard",
            "size_distribution": [
                {"size_pt": size, "count": count}
                for size, count in sorted_sizes[:10]
            ],
            "title_sizes": title_sizes[:3],
            "body_sizes": body_sizes[:3],
            "bold_ratio": round(bold_count / total_runs, 2) if total_runs > 0 else 0,
        }

    def _extract_layouts(self) -> Dict[str, Any]:
        """레이아웃 패턴 추출"""
        layout_patterns = []

        for i, slide in enumerate(self.prs.slides):
            shapes_info = []
            for shape in slide.shapes:
                info = {
                    "type": shape.shape_type.__class__.__name__ if hasattr(shape, 'shape_type') else "unknown",
                    "left": round(shape.left / 914400, 2) if shape.left else 0,
                    "top": round(shape.top / 914400, 2) if shape.top else 0,
                    "width": round(shape.width / 914400, 2) if shape.width else 0,
                    "height": round(shape.height / 914400, 2) if shape.height else 0,
                    "has_text": shape.has_text_frame,
                }
                shapes_info.append(info)

            # 패턴 유형 분류
            pattern = self._classify_layout(shapes_info)
            layout_patterns.append({
                "slide_index": i + 1,
                "shape_count": len(shapes_info),
                "pattern": pattern,
            })

        pattern_counter = Counter(p["pattern"] for p in layout_patterns)

        return {
            "slides": layout_patterns,
            "pattern_distribution": dict(pattern_counter.most_common()),
            "avg_shapes_per_slide": round(
                sum(p["shape_count"] for p in layout_patterns) / len(layout_patterns), 1
            ) if layout_patterns else 0,
        }

    def _extract_structure(self) -> Dict[str, Any]:
        """슬라이드 구조 분석 (섹션 구분, 유형 추정)"""
        structure = []

        for i, slide in enumerate(self.prs.slides):
            texts = []
            max_font_size = 0
            has_image = False
            has_table = False
            bg_is_dark = False

            # 배경색 체크
            bg = self._get_slide_bg_color(slide)
            if bg and self._brightness(bg) < 80:
                bg_is_dark = True

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text[:100])
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt > max_font_size:
                                max_font_size = run.font.size.pt

                if hasattr(shape, 'image'):
                    has_image = True
                if shape.has_table:
                    has_table = True

            # 슬라이드 유형 추정
            slide_type = self._guess_slide_type(
                texts, max_font_size, has_image, has_table, bg_is_dark, i
            )

            structure.append({
                "slide_index": i + 1,
                "type": slide_type,
                "title": texts[0] if texts else "",
                "text_count": len(texts),
                "max_font_size": max_font_size,
                "has_image": has_image,
                "has_table": has_table,
                "bg_dark": bg_is_dark,
            })

        type_counter = Counter(s["type"] for s in structure)

        return {
            "slides": structure,
            "type_distribution": dict(type_counter.most_common()),
            "section_dividers": [s for s in structure if s["type"] == "section_divider"],
        }

    # ─── 테마 변환 ─────────────────────────────────────────────

    def to_slide_kit_theme(self) -> Dict[str, Tuple[int, int, int]]:
        """분석 결과를 slide_kit apply_theme() 호환 형식으로 변환

        Returns:
            {"primary": (r,g,b), "secondary": (r,g,b), ...} 형식의 딕셔너리
        """
        analysis = self.analyze()
        colors = analysis["colors"]

        # primary: 가장 많이 사용된 어두운 색상
        primary = (0, 44, 95)  # 기본값
        if colors["primary_candidates"]:
            primary = tuple(colors["primary_candidates"][0]["rgb"])

        # secondary: 가장 많이 사용된 밝은 색상
        secondary = (0, 170, 210)  # 기본값
        if colors["secondary_candidates"]:
            secondary = tuple(colors["secondary_candidates"][0]["rgb"])

        # teal: secondary와 다른 밝은 색상 또는 기본값
        teal = (0, 161, 156)
        if len(colors["secondary_candidates"]) >= 2:
            teal = tuple(colors["secondary_candidates"][1]["rgb"])

        # accent: primary/secondary와 색상 대비가 큰 것
        accent = (230, 51, 18)
        for c in colors["primary_candidates"][1:] + colors["secondary_candidates"][1:]:
            rgb = tuple(c["rgb"])
            if self._color_distance(rgb, primary) > 150 and self._color_distance(rgb, secondary) > 100:
                accent = rgb
                break

        # dark/light 배경색
        dark = (33, 33, 33)
        light = (245, 245, 245)
        for bg in colors["background_colors"]:
            rgb = tuple(bg["rgb"])
            if self._brightness(rgb) < 50:
                dark = rgb
            elif self._brightness(rgb) > 200 and rgb != (255, 255, 255):
                light = rgb

        theme = {
            "primary": primary,
            "secondary": secondary,
            "teal": teal,
            "accent": accent,
            "dark": dark,
            "light": light,
        }

        logger.info(
            f"테마 변환 완료: primary={self._rgb_to_hex(primary)}, "
            f"secondary={self._rgb_to_hex(secondary)}"
        )

        return theme

    def to_design_profile(self) -> Dict[str, Any]:
        """분석 결과를 디자인 프로파일 형식으로 반환

        CLAUDE.md 및 생성 스크립트에서 참조할 수 있는 구조화된 정보
        """
        analysis = self.analyze()
        theme = self.to_slide_kit_theme()
        fonts = analysis["fonts"]

        return {
            "theme": theme,
            "theme_hex": {k: self._rgb_to_hex(v) for k, v in theme.items()},
            "primary_font": fonts["primary_font"],
            "title_sizes": fonts["title_sizes"],
            "body_sizes": fonts["body_sizes"],
            "slide_count": analysis["slide_count"],
            "structure_summary": analysis["structure"]["type_distribution"],
            "layout_summary": analysis["layouts"]["pattern_distribution"],
        }

    def save_analysis(self, output_path: Path) -> Path:
        """분석 결과를 JSON 파일로 저장"""
        analysis = self.analyze()
        analysis["slide_kit_theme"] = self.to_slide_kit_theme()
        analysis["design_profile"] = self.to_design_profile()

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # JSON 직렬화 (tuple → list)
        def convert(obj):
            if isinstance(obj, tuple):
                return list(obj)
            if isinstance(obj, RGBColor):
                return [obj[0], obj[1], obj[2]]
            return obj

        json_str = json.dumps(analysis, default=convert, ensure_ascii=False, indent=2)
        output_path.write_text(json_str, encoding="utf-8")

        logger.info(f"분석 결과 저장: {output_path}")
        return output_path

    # ─── 유틸리티 메서드 ──────────────────────────────────────

    def _get_slide_bg_color(self, slide) -> Optional[Tuple[int, int, int]]:
        """슬라이드 배경색 추출"""
        try:
            bg = slide.background
            if bg.fill.type is not None:
                if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color and bg.fill.fore_color.rgb:
                    rgb = bg.fill.fore_color.rgb
                    return (rgb[0], rgb[1], rgb[2])
        except Exception:
            pass
        return None

    def _get_shape_fill_color(self, shape) -> Optional[Tuple[int, int, int]]:
        """도형 채우기 색상 추출"""
        try:
            fill = shape.fill
            if fill.type is not None and hasattr(fill, 'fore_color'):
                if fill.fore_color and fill.fore_color.rgb:
                    rgb = fill.fore_color.rgb
                    return (rgb[0], rgb[1], rgb[2])
        except Exception:
            pass
        return None

    @staticmethod
    def _is_neutral(color: Tuple[int, int, int], threshold: int = 30) -> bool:
        """무채색(흰, 검, 회색) 여부 판단"""
        r, g, b = color
        # 흰색 계열
        if min(r, g, b) > 220:
            return True
        # 검정 계열
        if max(r, g, b) < 35:
            return True
        # 회색 계열 (R≈G≈B)
        if max(r, g, b) - min(r, g, b) < threshold:
            return True
        return False

    @staticmethod
    def _brightness(color: Tuple[int, int, int]) -> float:
        """색상 밝기 계산 (0~255)"""
        r, g, b = color
        return 0.299 * r + 0.587 * g + 0.114 * b

    @staticmethod
    def _rgb_to_hex(color: Tuple[int, int, int]) -> str:
        """RGB → #HEX 변환"""
        return "#{:02X}{:02X}{:02X}".format(*color)

    @staticmethod
    def _color_distance(c1: Tuple[int, int, int], c2: Tuple[int, int, int]) -> float:
        """두 색상 간 유클리드 거리"""
        return sum((a - b) ** 2 for a, b in zip(c1, c2)) ** 0.5

    def _classify_layout(self, shapes: List[Dict]) -> str:
        """도형 배치 기반 레이아웃 패턴 분류"""
        if not shapes:
            return "blank"

        text_shapes = [s for s in shapes if s["has_text"]]
        count = len(shapes)

        # 도형 1~2개: 단순 텍스트 (표지, 구분자)
        if count <= 2:
            return "minimal"

        # 도형이 가로로 나란히 배치 (멀티 컬럼)
        tops = sorted(set(s["top"] for s in shapes))
        if len(text_shapes) >= 3:
            # 같은 Y좌표에 3개 이상 → 그리드/카드
            top_groups = Counter(round(s["top"], 1) for s in text_shapes)
            max_same_row = max(top_groups.values()) if top_groups else 0
            if max_same_row >= 3:
                return "grid"
            elif max_same_row == 2:
                return "two_column"

        # 세로 스택
        if len(tops) >= 4:
            return "stacked"

        return "mixed"

    def _guess_slide_type(
        self,
        texts: List[str],
        max_font_size: float,
        has_image: bool,
        has_table: bool,
        bg_is_dark: bool,
        index: int,
    ) -> str:
        """슬라이드 유형 추정"""
        # 첫 슬라이드 = 표지
        if index == 0:
            return "cover"

        # 마지막 슬라이드 = 클로징
        if index == len(self.prs.slides) - 1:
            return "closing"

        # 다크 배경 + 큰 텍스트 = 섹션 구분자
        if bg_is_dark and max_font_size >= 36:
            return "section_divider"

        # 큰 폰트 + 텍스트 적음 = 강조/키메시지
        if max_font_size >= 36 and len(texts) <= 3:
            return "key_message"

        # 테이블 포함
        if has_table:
            return "table"

        # 이미지 포함
        if has_image:
            return "visual"

        # 텍스트 많음 = 콘텐츠
        if len(texts) >= 4:
            return "content"

        return "content"


def analyze_reference(pptx_path: Union[str, Path]) -> Dict[str, Any]:
    """레퍼런스 PPTX 파일 분석 (편의 함수)

    Args:
        pptx_path: 레퍼런스 PPTX 파일 경로

    Returns:
        디자인 프로파일 딕셔너리
    """
    analyzer = ReferenceAnalyzer(pptx_path)
    return analyzer.to_design_profile()


def analyze_and_apply_theme(pptx_path: Union[str, Path]) -> Dict[str, Any]:
    """레퍼런스 PPTX 분석 후 slide_kit 테마 자동 적용

    Args:
        pptx_path: 레퍼런스 PPTX 파일 경로

    Returns:
        적용된 테마 딕셔너리
    """
    analyzer = ReferenceAnalyzer(pptx_path)
    theme = analyzer.to_slide_kit_theme()

    # slide_kit의 C 딕셔너리에 직접 적용
    try:
        from ..generators.slide_kit import C, RGBColor as _RGB
        for key, rgb in theme.items():
            if key in C:
                C[key] = _RGB(*rgb)
        logger.info("slide_kit 컬러 팔레트에 레퍼런스 테마 적용 완료")
    except ImportError:
        logger.warning("slide_kit import 실패 — 테마 미적용")

    return theme
